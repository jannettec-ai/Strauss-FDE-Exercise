"""
build_deck.py

Generates the FDE presentation deck as a .pptx file.
Run: python scripts/build_deck.py
Output: docs/strauss_fde_deck.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand colours ─────────────────────────────────────────────────────────────

RED    = RGBColor(0xC8, 0x10, 0x2E)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0xF0, 0xF0, 0xF0)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DRED   = RGBColor(0x96, 0x0C, 0x22)

W = Inches(13.33)   # slide width  (16:9)
H = Inches(7.5)     # slide height

OUT = Path(__file__).parent.parent / "docs" / "strauss_fde_deck.pptx"

# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs: Presentation):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def txb(slide, left, top, width, height):
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )


def para(tf, text, size, bold=False, color=DARK, align=PP_ALIGN.LEFT,
         space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.text = text
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    p.space_before = Pt(space_before)
    p.space_after  = Pt(space_after)
    return p


def red_bar(slide, top=0.0, height=0.12):
    """Full-width red bar — used as header accent."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(top), W, Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    return bar


def red_left_bar(slide, top=1.0, height=5.5):
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(top), Inches(0.08), Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()


def box(slide, left, top, width, height, fill=LGRAY):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    return shape


def slide_header(slide, title, subtitle=None):
    red_bar(slide, top=0.0, height=0.1)
    tb = txb(slide, 0.2, 0.1, 12.9, 0.75)
    tf = tb.text_frame
    tf.word_wrap = False
    para(tf, title, 28, bold=True, color=DARK)
    if subtitle:
        tb2 = txb(slide, 0.2, 0.75, 12.9, 0.4)
        para(tb2.text_frame, subtitle, 13, color=GRAY)
    red_bar(slide, top=H.inches - 0.22, height=0.22)
    # footer
    ft = txb(slide, 0.2, H.inches - 0.22, 12.9, 0.22)
    para(ft.text_frame, "Strauss Procurement · FDE Exercise · Confidential",
         9, color=WHITE, align=PP_ALIGN.LEFT)


def bullets(slide, left, top, width, height, items, size=15, indent=False):
    tb = txb(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            p.text = item
            first = False
        else:
            p = tf.add_paragraph()
            p.text = item
        p.alignment = PP_ALIGN.LEFT
        r = p.runs[0] if p.runs else p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = DARK
        r.font.name = "Calibri"
        p.space_before = Pt(4 if indent else 6)
        p.level = 1 if indent else 0


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs):
    """1 — Title slide."""
    sl = blank(prs)

    # Full red left panel
    panel = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(5.5), H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RED
    panel.line.fill.background()

    # White logo text on red
    tb = txb(sl, 0.35, 0.4, 4.8, 1.0)
    para(tb.text_frame, "STRAUSS GROUP", 18, bold=True, color=WHITE)

    tb2 = txb(sl, 0.35, 1.35, 4.8, 3.5)
    tf = tb2.text_frame
    tf.word_wrap = True
    para(tf, "Negotiation Prep Intelligence", 32, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.text = "FDE Home Exercise"
    p2.runs[0].font.size = Pt(18)
    p2.runs[0].font.color.rgb = RGBColor(0xF5, 0xC0, 0xC8)
    p2.runs[0].font.name = "Calibri"
    p2.space_before = Pt(12)

    tb3 = txb(sl, 0.35, 6.4, 4.8, 0.8)
    tf3 = tb3.text_frame
    para(tf3, "Jannette · June 2026", 13, color=RGBColor(0xF5, 0xC0, 0xC8))

    # Right panel — tagline
    tb4 = txb(sl, 6.0, 2.2, 6.8, 3.5)
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    para(tf4,
         '"I want my team spending their time on judgment\nand relationships.\nNot on logistics."',
         20, color=GRAY, align=PP_ALIGN.LEFT)
    tb5 = txb(sl, 6.0, 5.2, 6.8, 0.6)
    para(tb5.text_frame, "— Strauss Procurement Manager", 13, color=GRAY)


def slide_problem(prs):
    """2 — The problem."""
    sl = blank(prs)
    slide_header(sl, "The Problem", "What the procurement team is actually spending time on")

    # Three stat boxes
    stats = [
        ("45–90 min", "of manual prep\nper meeting"),
        ("10–20", "supplier meetings\nper month"),
        ("8–30 hrs", "of information work\nper month, per team"),
    ]
    for i, (num, label) in enumerate(stats):
        bx = box(sl, 0.35 + i * 4.3, 1.4, 3.9, 1.6)
        tb = txb(sl, 0.45 + i * 4.3, 1.45, 3.7, 0.8)
        para(tb.text_frame, num, 36, bold=True, color=RED, align=PP_ALIGN.CENTER)
        tb2 = txb(sl, 0.45 + i * 4.3, 2.15, 3.7, 0.7)
        para(tb2.text_frame, label, 13, color=GRAY, align=PP_ALIGN.CENTER)

    tb_desc = txb(sl, 0.35, 3.25, 12.6, 0.6)
    para(tb_desc.text_frame,
         "The Procurement Manager called it 'information logistics' \u2014 work that stands between the team and the judgment call they're paid to make.",
         14, color=DARK)

    # Four failure modes
    tb_h = txb(sl, 0.35, 3.95, 12.6, 0.4)
    para(tb_h.text_frame, "Four ways the manual process fails consistently:", 13, bold=True, color=RED)

    failures = [
        "❶  Wrong contract version pulled → missed penalty clause before negotiation",
        "❷  Open dispute threads go cold → unresolved delivery or quality issues forgotten by next meeting",
        "❸  Benchmark data stale → manager walks into cocoa or sugar negotiation without knowing where market moved",
        "❹  Expired contract not flagged → team operates on informal terms for months before anyone notices",
    ]
    bullets(sl, 0.35, 4.4, 12.6, 2.8, failures, size=13)


def slide_what_built(prs):
    """3 — What I built."""
    sl = blank(prs)
    slide_header(sl, "What I Built", "One tool. One job. One page per meeting.")

    # Big pitch box
    bx = box(sl, 0.35, 1.3, 12.6, 1.5, fill=RGBColor(0xFB, 0xEB, 0xED))
    tb = txb(sl, 0.55, 1.4, 12.2, 1.3)
    tf = tb.text_frame
    tf.word_wrap = True
    para(tf,
         "An AI-generated one-page briefing, produced automatically before every supplier meeting, "
         "that pulls together email history, contract terms, and the commodity price benchmark "
         "into a single view — with every claim traceable to its source.",
         16, color=DARK)

    tb_url = txb(sl, 0.35, 3.0, 12.6, 0.45)
    para(tb_url.text_frame, "🔗  strauss-procurement.streamlit.app", 13, color=RED)

    # Three columns: what it does
    cols = [
        ("📧  Emails", "Summarises 6 months of supplier thread. Flags open disputes, unanswered threads, and price quotes."),
        ("📄  Contracts", "Extracts payment terms, volume commitments, penalty clauses, and renewal dates from the active contract version."),
        ("📈  Prices", "Adds the live commodity benchmark (cocoa, sugar, dairy) so the manager knows where the market sits before walking in."),
    ]
    for i, (title, body) in enumerate(cols):
        bx2 = box(sl, 0.35 + i * 4.3, 3.6, 4.0, 3.2)
        tb3 = txb(sl, 0.5 + i * 4.3, 3.7, 3.7, 0.5)
        para(tb3.text_frame, title, 14, bold=True, color=RED)
        tb4 = txb(sl, 0.5 + i * 4.3, 4.2, 3.7, 2.4)
        tf4 = tb4.text_frame
        tf4.word_wrap = True
        para(tf4, body, 13, color=DARK)


def slide_architecture(prs):
    """4 — Architecture."""
    sl = blank(prs)
    slide_header(sl, "How It Works", "Prototype vs. production — a deliberate distinction")

    # Left: prototype
    box(sl, 0.35, 1.3, 5.8, 5.8, fill=LGRAY)
    tb = txb(sl, 0.55, 1.35, 5.4, 0.45)
    para(tb.text_frame, "PROTOTYPE  (what you're seeing today)", 11, bold=True, color=GRAY)

    proto_lines = [
        "Emails → JSON files (synthetic, ~111 emails)",
        "Contracts → Markdown files (10 contracts)",
        "Prices → Static CSV (yfinance, one-time pull)",
        "",
        "↓  extraction.py",
        "One Claude API call per supplier",
        "Extracts: open issues · price quotes",
        "contract terms · commodity commentary",
        "",
        "↓  packet_generator.py",
        "Assembles packet per meeting ID",
        "Computes KPIs locally (no API)",
        "",
        "↓  Streamlit UI",
        "Home · Suppliers · Meeting Prep",
    ]
    bullets(sl, 0.55, 1.85, 5.4, 5.0, proto_lines, size=12)

    # Arrow
    tb_arr = txb(sl, 6.25, 3.8, 0.8, 0.6)
    para(tb_arr.text_frame, "→", 28, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # Right: production
    box(sl, 7.1, 1.3, 5.85, 5.8, fill=RGBColor(0xFB, 0xEB, 0xED))
    tb2 = txb(sl, 7.3, 1.35, 5.5, 0.45)
    para(tb2.text_frame, "PRODUCTION  (Phase 2 architecture)", 11, bold=True, color=RED)

    prod_lines = [
        "Emails → Microsoft Graph API (Outlook)",
        "Contracts → SharePoint + OCR pipeline",
        "Prices → Bloomberg / ICE live feed",
        "Connector layer: MCP per source",
        "",
        "↓  Storage split",
        "Structured store (BigQuery): exact facts",
        "Vector store: semantic email search",
        "",
        "↓  Packet Generator",
        "Same logic, real-time data",
        "Triggered by calendar event",
        "",
        "↓  Streamlit / Teams tab",
        "Auto-delivered before each meeting",
    ]
    bullets(sl, 7.3, 1.85, 5.5, 5.0, prod_lines, size=12)


def slide_demo_home(prs):
    """5 — Demo: Home dashboard."""
    sl = blank(prs)
    slide_header(sl, "Demo: Home Dashboard", "Instant situational awareness — no API calls, loads in <1s")

    items = [
        "📊  Stat tiles — upcoming meetings, meetings this week, suppliers at risk, contract issues",
        "📅  Full meeting calendar — supplier, date, agenda note, geography, urgency badge",
        "🏭  Supplier health grid — relationship status at a glance across all 8 suppliers",
        "📧  Latest discussions — most recent email per supplier with inline preview",
    ]
    bullets(sl, 0.35, 1.35, 7.5, 3.0, items, size=15)

    # Key design call-out
    bx = box(sl, 0.35, 4.5, 7.5, 2.65, fill=RGBColor(0xFB, 0xEB, 0xED))
    tb = txb(sl, 0.55, 4.6, 7.1, 0.4)
    para(tb.text_frame, "Design decision", 12, bold=True, color=RED)
    tb2 = txb(sl, 0.55, 5.0, 7.1, 2.0)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    para(tf2,
         "The home page makes zero API calls. All analytics are computed locally "
         "from email and contract files. A Procurement Manager can load the dashboard "
         "at any time — including 5 minutes before a meeting — without cost or latency.",
         13, color=DARK)

    # Right: screenshot placeholder
    bx2 = box(sl, 8.1, 1.35, 4.85, 5.8, fill=RGBColor(0xE8, 0xE8, 0xE8))
    tb3 = txb(sl, 8.1, 3.5, 4.85, 0.8)
    para(tb3.text_frame, "[ Live demo ]", 16, color=GRAY, align=PP_ALIGN.CENTER)
    tb4 = txb(sl, 8.1, 4.1, 4.85, 0.5)
    para(tb4.text_frame, "strauss-procurement.streamlit.app", 10, color=GRAY, align=PP_ALIGN.CENTER)


def slide_demo_supplier(prs):
    """6 — Demo: Supplier profile."""
    sl = blank(prs)
    slide_header(sl, "Demo: Supplier Profile", "Relationship health analytics — searchable directory, per-supplier depth")

    items = [
        "🔍  Searchable directory — filter by name, category, or geography",
        "🟢🟡🔴  Relationship health score — composite of response time trend + communication volume",
        "📊  Email activity chart — emails per month, last 12 months",
        "📈  Responsiveness trend — early-half vs recent-half avg response time (getting slower = warning)",
        "📉  Retention signal — volume trend, days since last contact, interpretation in plain English",
        "📄  Latest discussion preview — most recent email thread with inline read",
        "🔗  One-click shortcut to generate the meeting prep packet for this supplier",
    ]
    bullets(sl, 0.35, 1.35, 8.5, 5.5, items, size=14)

    bx = box(sl, 9.1, 1.35, 3.85, 5.8, fill=RGBColor(0xE8, 0xE8, 0xE8))
    tb = txb(sl, 9.1, 3.5, 3.85, 0.8)
    para(tb.text_frame, "[ Live demo ]", 16, color=GRAY, align=PP_ALIGN.CENTER)


def slide_demo_packet(prs):
    """7 — Demo: Meeting prep packet."""
    sl = blank(prs)
    slide_header(sl, "Demo: Meeting Prep Packet", "The core product — one click, one page, everything the manager needs")

    left_items = [
        "⚡  Select meeting → Generate Packet",
        "One Claude API call per supplier",
        "Returns in ~10–15 seconds",
        "",
        "Packet contains:",
        "  • Heads-up alert (most critical issue)",
        "  • 4 KPI tiles (response time, open issues,",
        "    price delta, days to renewal)",
        "  • Open issues list (colour-coded by severity)",
        "  • Email thread summary",
        "  • Contract at a glance",
        "  • Pricing section with delta calculation",
        "  • Commodity benchmark (6-month change)",
    ]
    bullets(sl, 0.35, 1.35, 6.2, 5.8, left_items, size=13)

    # Galil call-out
    bx = box(sl, 6.8, 1.35, 6.15, 5.8, fill=RGBColor(0xFB, 0xEB, 0xED))
    tb = txb(sl, 7.0, 1.45, 5.8, 0.45)
    para(tb.text_frame, "Demo case: Galil Dairy Suppliers", 13, bold=True, color=RED)

    galil_items = [
        "🔴  Price discrepancy flagged",
        "     Supplier quoted ₪15.50/kg in email",
        "     Contract Section 4.2 caps at ₪14.20/kg",
        "     No written amendment exists",
        "     Delta: +₪1.30/kg (+9.2%)",
        "",
        "🟡  Unanswered thread flagged",
        "     Yael Shapira requested alignment",
        "     David Cohen proposed call — no",
        "     documented outcome",
        "",
        "📎  Every fact has a 'Show source'",
        "     expander pointing to the exact",
        "     email or contract section",
    ]
    bullets(sl, 7.0, 1.95, 5.8, 5.0, galil_items, size=12)


def slide_human_loop(prs):
    """8 — Human in the loop."""
    sl = blank(prs)
    slide_header(sl, "The Human in the Loop", "What stops this from being trusted blindly")

    mechanisms = [
        (
            "📎  Source attribution on every field",
            "Every number, date, and clause cites exactly where it came from "
            "(email date + sender, or contract section). "
            "One click to verify. Nothing is presented as settled ground truth."
        ),
        (
            "⚠️  Explicit conflict flagging",
            "When email and contract disagree, the packet surfaces both and flags the conflict "
            "— it never silently picks one. Galil price discrepancy and Siam Sugar unit ambiguity "
            "are live examples in the demo."
        ),
        (
            "🚩  Correction tracking built in",
            "The manager can flag any field as wrong directly in the packet view. "
            "Corrections are saved to a log (correction_log.csv) — this is the "
            "extraction quality feedback loop, per metrics.md §5."
        ),
        (
            "🚫  Tool never acts",
            "No emails sent, no contracts updated, no decisions taken. "
            "Every action stays with the human. "
            "The tool removes the assembly work in front of the judgment call — it does not replace the judgment call."
        ),
    ]
    for i, (title, body) in enumerate(mechanisms):
        top = 1.35 + i * 1.45
        bx = box(sl, 0.35, top, 12.6, 1.3)
        tb1 = txb(sl, 0.55, top + 0.08, 12.2, 0.4)
        para(tb1.text_frame, title, 14, bold=True, color=RED)
        tb2 = txb(sl, 0.55, top + 0.48, 12.2, 0.75)
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        para(tf2, body, 13, color=DARK)


def slide_kpis(prs):
    """12 — KPIs with FDE tracking column."""
    sl = blank(prs)
    slide_header(sl, "Five KPIs — Measurable in 90 Days",
                 "Defined in metrics.md · Captured per packet in SQLite · Visible in FDE dashboard")

    # Column headers
    for left, width, label in [
        (0.2, 0.5,  ""),
        (0.8, 4.5,  "Definition & formula"),
        (5.4, 3.5,  "How it is computed"),
        (9.0, 2.0,  "Live example"),
        (11.1, 2.0, "FDE tracking"),
    ]:
        bx = box(sl, left, 1.28, width, 0.38, fill=RED)
        if label:
            tb = txb(sl, left + 0.05, 1.3, width - 0.08, 0.34)
            para(tb.text_frame, label, 10, bold=True, color=WHITE)

    kpis = [
        (
            "§1", "Supplier Response Time",
            "Avg calendar days between each Strauss outbound email and the first supplier reply, across the 6-month window.",
            "Computed locally from email date fields. No API call needed. Pairs matched chronologically per thread.",
            "Galil Dairy: 11.2d avg",
            "kpi_response_days stored per packet. Trend visible in per-supplier breakdown.",
        ),
        (
            "§2", "Open Issues Count",
            "Number of distinct unresolved threads or flagged items per supplier as of the most recent email date.",
            "Claude API extraction pass. Each issue has type, description, source_ref, and status=open.",
            "Galil: 2 open (price discrepancy + unanswered thread)",
            "kpi_open_issues stored per packet. Avg across all packets shown in KPI tile.",
        ),
        (
            "§3", "Price vs. Contract Delta",
            "Latest email price quote minus contract base price. Absolute and percentage. Alert if delta > 0 without written amendment.",
            "Computed locally. latest_price_quoted.value from Claude vs contract_base_price_numeric from contract.",
            "Galil: +9.2% (+NIS 1.30/kg). Triggers red alert.",
            "kpi_price_delta_pct stored per packet. Tracks which suppliers are consistently over contract.",
        ),
        (
            "§4", "Days to Renewal",
            "Calendar days from today to contract renewal date. Red <30d, amber 30-90d. Auto-renewal shows notice deadline.",
            "Computed locally from contract renewal_date field vs system date. No API call.",
            "Lowlands Dairy: no active contract flagged immediately.",
            "kpi_days_to_renewal stored per packet. Surfaces suppliers approaching critical renewal window.",
        ),
        (
            "§5", "Override / Correction Rate",
            "% of AI-extracted fields flagged as wrong per packet by the manager. High rate = extraction needs tuning.",
            "UI checkboxes in packet view. Saved to correction_log.csv when manager clicks Save. Joined to SQLite for aggregate view.",
            "Target: <10% by day 90. Currently tracked per session.",
            "kpi_correction_rate_pct written when manager saves corrections. Identifies which fields and suppliers have highest error rate.",
        ),
    ]

    row_fills = [LGRAY, WHITE, LGRAY, WHITE, LGRAY]
    for i, (tag, title, definition, computation, example, fde) in enumerate(kpis):
        top = 1.72 + i * 1.13
        fill = row_fills[i]
        for left, width in [(0.2, 0.5), (0.8, 4.5), (5.4, 3.5), (9.0, 2.0), (11.1, 2.0)]:
            box(sl, left, top, width, 1.08, fill=fill)

        # Tag
        tag_b = sl.shapes.add_shape(1, Inches(0.22), Inches(top + 0.08), Inches(0.44), Inches(0.28))
        tag_b.fill.solid(); tag_b.fill.fore_color.rgb = RED; tag_b.line.fill.background()
        tb_tag = txb(sl, 0.23, top + 0.09, 0.42, 0.25)
        para(tb_tag.text_frame, tag, 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Title
        tb_t = txb(sl, 0.23, top + 0.42, 0.44, 0.55)
        tf_t = tb_t.text_frame; tf_t.word_wrap = True
        para(tf_t, title, 8, bold=True, color=DARK)

        for left, width, text in [
            (0.85, 4.35, definition),
            (5.45, 3.4,  computation),
            (9.05, 1.9,  example),
            (11.15, 1.9, fde),
        ]:
            tb = txb(sl, left, top + 0.06, width, 0.95)
            tf = tb.text_frame; tf.word_wrap = True
            para(tf, text, 9, color=DARK if text != fde else GRAY)


def slide_production(prs):
    """New — Production architecture."""
    sl = blank(prs)
    slide_header(sl, "Production Architecture",
                 "How the prototype becomes a production system — and what each layer would look like")

    # Three layers: Data Sources, Processing, Outputs + Monitoring
    layers = [
        (
            "DATA SOURCES",
            RED,
            [
                "Emails",
                "Microsoft Graph API",
                "Outlook / Exchange",
                "Shared procurement mailbox",
                "Push subscription on new mail",
                "",
                "Contracts",
                "SharePoint connector",
                "PDF extraction + OCR pipeline",
                "Version detection (active vs draft)",
                "",
                "Prices",
                "Bloomberg / ICE / Refinitiv",
                "Structured time-series feed",
                "Normalized to common schema",
            ]
        ),
        (
            "PROCESSING LAYER",
            DRED,
            [
                "MCP connector layer",
                "Standardised tool interface",
                "Swappable backend per source",
                "",
                "Storage split",
                "Structured store (BigQuery):",
                "  exact facts — contract terms,",
                "  KPI values, price points",
                "",
                "Vector store (Vertex AI / Pinecone):",
                "  chunked email bodies,",
                "  contract clause language,",
                "  semantic search at scale",
                "",
                "Packet Generator",
                "Same logic as prototype",
                "Triggered by calendar event 24h",
                "before each meeting",
            ]
        ),
        (
            "OUTPUTS + MONITORING",
            GRAY,
            [
                "Procurement Manager UI",
                "Streamlit or Teams tab",
                "Packet auto-delivered before meeting",
                "Source attribution + conflict flags",
                "",
                "FDE Metrics Pipeline",
                "Generation events: Pub/Sub",
                "  → BigQuery events table",
                "KPI values per packet stored",
                "Correction log aggregated",
                "Looker / Data Studio dashboard",
                "for ongoing monitoring",
                "",
                "Prompt tuning cycle",
                "Monthly: correction_log patterns",
                "  → extraction prompt update",
                "  → re-run accuracy benchmark",
            ]
        ),
    ]

    for i, (title, color, lines) in enumerate(layers):
        left = 0.2 + i * 4.37
        bar = sl.shapes.add_shape(1, Inches(left), Inches(1.28), Inches(4.2), Inches(0.38))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
        tb_h = txb(sl, left + 0.1, 1.3, 4.0, 0.34)
        para(tb_h.text_frame, title, 11, bold=True, color=WHITE)
        bx = box(sl, left, 1.68, 4.2, 5.55)
        bullets(sl, left + 0.15, 1.75, 3.9, 5.4, lines, size=11)

    # Bottom note
    tb_note = txb(sl, 0.2, 7.1, 12.9, 0.3)
    para(tb_note.text_frame,
         "Prototype delta: flat JSON files instead of connectors; local extraction instead of vector store; "
         "SQLite instead of BigQuery. Every simplification is documented and justified.",
         10, color=GRAY)


def slide_rollout(prs):
    """10 — Rollout plan."""
    sl = blank(prs)
    slide_header(sl, "Rollout Plan", "Three phases — trust first, scale second")

    phases = [
        (
            "Phase 1  ·  Weeks 1–4",
            "Controlled pilot",
            [
                "2–3 upcoming meetings with one analyst as primary user",
                "Analyst runs the tool, reviews output, flags corrections",
                "Track correction rate weekly — target <20% by week 4",
                "One structured debrief after each meeting: what did the packet get right?",
                "Success gate: packet catches one real issue the analyst would have missed",
            ],
            RED,
        ),
        (
            "Phase 2  ·  Weeks 5–8",
            "Expanded pilot",
            [
                "All 8 suppliers in scope; full procurement team using it",
                "IT review: security, data access, API key governance",
                "Integrate with calendar to auto-trigger packet generation 24h before each meeting",
                "Add live data sources: Graph API for emails, SharePoint for contracts",
                "Success gate: team uses packet without being reminded; correction rate <10%",
            ],
            DRED,
        ),
        (
            "Phase 3  ·  Weeks 9–12",
            "Handover",
            [
                "Ownership transferred to internal champion (Procurement Manager)",
                "ROI case presented to Finance/CFO with 90-day actuals",
                "Prompt tuning backlog based on correction_log.csv patterns",
                "Roadmap: supplier segmentation, auto-renewal alerts, negotiation history",
            ],
            GRAY,
        ),
    ]

    for i, (phase, subtitle, items, color) in enumerate(phases):
        left = 0.35 + i * 4.32
        bx = box(sl, left, 1.3, 4.1, 5.85)
        bar = sl.shapes.add_shape(1, Inches(left), Inches(1.3), Inches(4.1), Inches(0.45))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        tb1 = txb(sl, left + 0.15, 1.32, 3.8, 0.38)
        para(tb1.text_frame, phase, 12, bold=True, color=WHITE)
        tb2 = txb(sl, left + 0.15, 1.82, 3.8, 0.38)
        para(tb2.text_frame, subtitle, 13, bold=True, color=color)
        bullets(sl, left + 0.15, 2.25, 3.8, 4.75, items, size=11)


def slide_change_management(prs):
    """11 — Change management."""
    sl = blank(prs)
    slide_header(sl, "Change Management", "Building it is 20% of the job — getting people to use it is the other 80%")

    rows = [
        (
            "Objection",
            "Who says it",
            "How to handle it",
        ),
        (
            '"The output is wrong."',
            "Procurement Analyst\n(highest adoption risk — feels threatened but won't say so)",
            "Involve them in Phase 1 QA from day one. Make the correction feature visibly theirs. "
            "Frame their role as reviewer, not recipient. The tool makes them look more prepared, not redundant.",
        ),
        (
            '"I don\'t trust AI to read\na contract."',
            "Legal / Compliance",
            "Every fact has a source citation. The design is 'surfaces for human review,' "
            "not 'interprets and acts.' Show them the correction log \u2014 the system has a visible trust mechanism.",
        ),
        (
            '"What does this cost\nand who owns it?"',
            "Finance / IT",
            "Phase 1 cost is API calls only (~$15/month). Nothing touches live systems. "
            "Static file architecture minimises security surface. "
            "ROI case ready for month 3 (next slide).",
        ),
    ]

    col_widths = [3.0, 3.2, 6.1]
    col_lefts  = [0.35, 3.45, 6.75]
    header_fill = RED
    row_fills   = [LGRAY, WHITE, LGRAY]

    for r_idx, row in enumerate(rows):
        is_header = r_idx == 0
        top = 1.3 + r_idx * 1.7
        height = 0.55 if is_header else 1.6
        for c_idx, (text, col_left, col_width) in enumerate(zip(row, col_lefts, col_widths)):
            fill = header_fill if is_header else row_fills[r_idx - 1]
            bx = box(sl, col_left, top, col_width, height, fill=fill)
            tb = txb(sl, col_left + 0.1, top + 0.08, col_width - 0.15, height - 0.12)
            tf = tb.text_frame
            tf.word_wrap = True
            fsize = 11 if is_header else (12 if c_idx == 0 else 11)
            fc = WHITE if is_header else (RED if c_idx == 0 else DARK)
            bold = is_header or c_idx == 0
            para(tf, text, fsize, bold=bold, color=fc)


def slide_roi(prs):
    """12 — ROI case."""
    sl = blank(prs)
    slide_header(sl, "The ROI Case", "Conservative anchor — lead with what's measurable, then the multiplier")

    # Left: the math
    box(sl, 0.35, 1.3, 7.0, 5.85)
    tb_h = txb(sl, 0.55, 1.38, 6.6, 0.4)
    para(tb_h.text_frame, "Labor cost avoidance — the conservative number", 13, bold=True, color=RED)

    math_lines = [
        "Assumption 1:  15 meetings/month (midpoint of stated 10–20)",
        "Assumption 2:  Prep time reduced from 60 min → 10 min per meeting",
        "                      (review packet + flag corrections)",
        "Time saved:      50 min × 15 meetings = 12.5 hrs/month",
        "",
        "Assumption 3:  Analyst fully-loaded cost: ₪200/hr",
        "                      [confirm with HR — this is the key input]",
        "",
        "Monthly saving:   12.5 hrs × ₪200 = ₪2,500/month",
        "Annual saving:    ₪2,500 × 12 = ₪30,000/year",
        "",
        "API cost:           ~15 meetings × ~$1/call = ~$15/month",
        "                      Negligible. Rounds to zero.",
        "",
        "Conservative 12-month ROI:   ₪30,000 in analyst time",
    ]
    bullets(sl, 0.55, 1.85, 6.6, 5.1, math_lines, size=12)

    # Right: multiplier
    bx2 = box(sl, 7.55, 1.3, 5.4, 2.8, fill=RGBColor(0xFB, 0xEB, 0xED))
    tb2 = txb(sl, 7.75, 1.38, 5.0, 0.4)
    para(tb2.text_frame, "The multiplier: one caught violation", 13, bold=True, color=RED)
    mul_lines = [
        "Galil Dairy (live in the demo):",
        "+₪1.30/kg quoted over contract cap",
        "22 tons/quarter × 4 quarters = 88 tons/year",
        "Unchallenged cost: ₪1.30 × 88,000 kg = ₪114,400",
        "",
        "The tool caught this in the prep packet.",
        "One meeting. One flagged clause.",
        "That alone covers 3+ years of the tool's cost.",
    ]
    bullets(sl, 7.75, 1.85, 5.0, 2.1, mul_lines, size=12)

    # Right: what to say to CFO
    bx3 = box(sl, 7.55, 4.3, 5.4, 2.85, fill=LGRAY)
    tb3 = txb(sl, 7.75, 4.38, 5.0, 0.4)
    para(tb3.text_frame, "What to tell the CFO at month 3", 13, bold=True, color=DARK)
    cfo_lines = [
        "Lead with: ₪30k/year in analyst time reclaimed.",
        "Don't lead with: speculative negotiation savings.",
        "Add: correction rate trending down (trust improving).",
        "Add: open issues caught before meetings (count them).",
        "The ₪114k violation multiplier is the kicker —",
        "present it as 'one example already found in pilot.'",
    ]
    bullets(sl, 7.75, 4.85, 5.0, 2.15, cfo_lines, size=12)


def slide_whats_next(prs):
    """13 — What's next."""
    sl = blank(prs)
    slide_header(sl, "What's Next", "What I'd prioritise in Phase 2 — and what I deliberately left out")

    left_items = [
        "Priority 1: Live email connector (Microsoft Graph API)",
        "    → Real supplier threads, no synthetic data",
        "    → Scheduled pull, triggered before each meeting",
        "",
        "Priority 2: SharePoint contract connector",
        "    → Always the active version, not a stale copy",
        "    → OCR pipeline for scanned PDFs",
        "",
        "Priority 3: Persist correction log → prompt tuning cycle",
        "    → correction_log.csv already captures this",
        "    → Monthly review: which fields/suppliers have high error rate?",
        "    → Update extraction prompt accordingly",
        "",
        "Priority 4: Supplier segmentation",
        "    → Flag high-risk suppliers based on relationship health score",
        "    → Auto-generate prep for every meeting, not just on-demand",
    ]
    bullets(sl, 0.35, 1.35, 7.5, 5.8, left_items, size=13)

    # Right: what I left out and why
    bx = box(sl, 8.1, 1.35, 4.85, 5.8, fill=LGRAY)
    tb = txb(sl, 8.3, 1.43, 4.5, 0.4)
    para(tb.text_frame, "Deliberately out of scope", 13, bold=True, color=RED)
    out_items = [
        "❌  Negotiation strategy advice",
        "     Human judgment stays human",
        "",
        "❌  Full contract lifecycle mgmt",
        "     Different problem, different tool",
        "",
        "❌  Live MCP connector in prototype",
        "     Demo-day network risk not worth it",
        "     — documented as production path",
        "",
        "❌  Vector store in prototype",
        "     Unjustified at 8 suppliers / 111 emails",
        "     — right call at real scale",
        "",
        "❌  Teams / Slack integration",
        "     Phase 3, after trust is established",
    ]
    bullets(sl, 8.3, 1.9, 4.5, 5.1, out_items, size=12)


def slide_stakeholders(prs):
    """3 - Stakeholder map."""
    sl = blank(prs)
    slide_header(sl, "Stakeholder Map", "Who uses this, who might resist, and my read on each")

    rows = [
        ("Stakeholder", "Role", "What they need", "My read"),
        ("Procurement\nManager",
         "Primary user\nInternal champion",
         "Speed + accuracy\nVisible conflict flags",
         "Strongest ally. Named the problem themselves.\nWill adopt fast if first packet catches something real.\nRisk: one factual error in demo = trust is hard to rebuild."),
        ("Procurement\nAnalyst",
         "Day-to-day user\nQA reviewer",
         "Reduction in manual work\nEnough trust to stake prep on it",
         "Highest adoption risk. Won't say they feel threatened.\nWill say 'the output is wrong' or 'it adds a review step.'\nCounter: involve them in Phase 1 QA. Make corrections visibly theirs."),
        ("IT / Systems",
         "Enabler for\nproduction path",
         "Security review\nData governance\nAPI key management",
         "Not a blocker for prototype (static files, nothing touches live systems).\nCritical gating stakeholder for Phase 2.\nNeeds early visibility, not a last-minute ask."),
        ("Legal /\nCompliance",
         "Risk owner",
         "AI output not treated\nas authoritative\nwithout human review",
         "Potentially cautious. Answer: every fact has a source citation.\nFrame as 'surfaces for review' not 'interprets contracts.'\nCorrection mechanism is the direct answer to their concern."),
        ("Finance / CFO",
         "Cost approver",
         "Measurable ROI\nin 90 days",
         "Not a day-one stakeholder. Bring the ROI case at month 3.\nLead with labor cost avoidance (conservative, auditable).\nNot speculative negotiation savings."),
    ]

    col_widths = [1.6, 2.0, 2.5, 6.1]
    col_lefts  = [0.2, 1.9, 4.0, 6.6]
    header_fill = RED
    row_fills = [LGRAY, WHITE, LGRAY, WHITE, LGRAY]

    for r_idx, row in enumerate(rows):
        is_header = r_idx == 0
        top = 1.28 + r_idx * 1.18
        height = 0.45 if is_header else 1.08
        for c_idx, (text, col_left, col_width) in enumerate(zip(row, col_lefts, col_widths)):
            fill = header_fill if is_header else row_fills[r_idx - 1]
            bx = box(sl, col_left, top, col_width, height, fill=fill)
            tb = txb(sl, col_left + 0.08, top + 0.05, col_width - 0.12, height - 0.08)
            tf = tb.text_frame
            tf.word_wrap = True
            fsize = 10 if is_header else (10 if c_idx == 3 else 11)
            fc = WHITE if is_header else (RED if c_idx == 0 else DARK)
            bold = is_header or c_idx == 0
            para(tf, text, fsize, bold=bold, color=fc)


def slide_week1(prs):
    """4 - Week 1 plan."""
    sl = blank(prs)
    slide_header(sl, "Week 1 Plan", "Five things before touching any technology")

    steps = [
        (
            "1  Shadow a prep session in real time",
            "Sit next to an analyst while they prepare for an actual upcoming meeting. "
            "Don't ask them to describe the process — watch them do it. "
            "Note which tabs are open, how long each step takes, where they get stuck, what they skip when short on time. "
            "The stated process and the actual process are almost always different."
        ),
        (
            "2  Interview the Procurement Manager separately",
            "Ask: what has gone wrong in a meeting because prep was incomplete? "
            "What did you wish you'd known walking in? "
            "What would you never trust a tool to get right? "
            "The second and third questions are as important as the first."
        ),
        (
            "3  Audit the actual data sources — not descriptions of them",
            "Request read-only access to: one supplier email thread in Outlook, one contract on SharePoint, the price benchmark spreadsheet. "
            "Look at real messiness: how many threads per supplier, what the file naming convention actually is, who owns the price spreadsheet. "
            "You cannot design a robust extraction layer without seeing the raw input."
        ),
        (
            "4  Map what a good packet looks like by working backwards",
            "Pick the three highest-stakes upcoming meetings. Ask the manager: if you had one page of perfect context for each, what would be on it? "
            "This surfaces the prioritisation the tool needs to make — what to include and, equally important, what to leave out."
        ),
        (
            "5  Identify the first real meeting to demo against",
            "Before building anything, agree on a specific upcoming meeting as the Phase 1 success test. "
            "Forces the stakeholder to commit to an evaluation moment. "
            "Means the first demo is not a toy example — it is something the manager actually cares about walking into."
        ),
    ]

    for i, (title, body) in enumerate(steps):
        top = 1.28 + i * 1.2
        bar = sl.shapes.add_shape(1, Inches(0.2), Inches(top + 0.08), Inches(0.06), Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RED
        bar.line.fill.background()
        tb1 = txb(sl, 0.4, top, 12.5, 0.38)
        para(tb1.text_frame, title, 13, bold=True, color=RED)
        tb2 = txb(sl, 0.4, top + 0.38, 12.5, 0.75)
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        para(tf2, body, 12, color=DARK)


def slide_rejected(prs):
    """6 - What I considered and rejected."""
    sl = blank(prs)
    slide_header(sl, "What I Considered and Rejected", "Every design decision has a tradeoff — these were mine")

    rows = [
        ("What I considered", "Rejected because", "Tradeoff accepted"),
        ("Commodity pricing\ndashboard",
         "Answers an adjacent problem, not the stated pain point.\nMost obvious build — every candidate goes here.",
         "Pricing folded in as one input to the packet, not the centrepiece."),
        ("General 'what was agreed'\ntracking layer",
         "A system-of-record problem, not an AI problem.\nBetter solved by a database than a model.",
         "Solved indirectly: packet pulls from accurate, current contract terms."),
        ("Live MCP connector\nfor pricing in prototype",
         "Live demo depending on third-party server responding correctly in the room.\nUnnecessary risk for zero added insight value.",
         "One-time static CSV pull. MCP documented as production path, not built into prototype."),
        ("Vector store +\nlive connectors",
         "Premature at this data volume: 8 suppliers, ~111 emails.\nReal value at hundreds of suppliers, thousands of emails.",
         "Local extraction pass against flat files. Vector store called out explicitly as production-scale answer."),
        ("Full contract lifecycle\nmanagement",
         "Too broad a problem for the time available.\nNot the time-sink the brief describes.",
         "Contract terms are read and surfaced, not managed end to end."),
    ]

    col_widths = [2.8, 5.0, 4.8]
    col_lefts  = [0.2, 3.1, 8.2]
    header_fill = RED
    row_fills = [LGRAY, WHITE, LGRAY, WHITE, LGRAY]

    for r_idx, row in enumerate(rows):
        is_header = r_idx == 0
        top = 1.28 + r_idx * 1.18
        height = 0.42 if is_header else 1.08
        for c_idx, (text, col_left, col_width) in enumerate(zip(row, col_lefts, col_widths)):
            fill = header_fill if is_header else row_fills[r_idx - 1]
            bx = box(sl, col_left, top, col_width, height, fill=fill)
            tb = txb(sl, col_left + 0.08, top + 0.05, col_width - 0.12, height - 0.08)
            tf = tb.text_frame
            tf.word_wrap = True
            fsize = 11 if is_header else 11
            fc = WHITE if is_header else (RED if c_idx == 0 else DARK)
            bold = is_header or c_idx == 0
            para(tf, text, fsize, bold=bold, color=fc)


def slide_risks(prs):
    """14 - Risks and challenges."""
    sl = blank(prs)
    slide_header(sl, "Risks and Challenges", "Technical and organisational — what could go wrong and how I would handle it")

    risks = [
        (
            "Extraction accuracy drift",
            "TECHNICAL",
            "Claude extraction returns wrong contract terms or misses a dispute. Manager walks in with bad data.",
            "Source attribution on every field makes errors visible. Correction log tracks error rate by field and supplier. "
            "Monthly prompt tuning cycle against correction_log.csv. Phase 1 KPI: correction rate <10% by week 4."
        ),
        (
            "Losing the internal champion",
            "ORGANISATIONAL",
            "Procurement Manager moves roles or loses interest. Tool loses its sponsor before Phase 2.",
            "Involve the analyst as a co-owner from day one — they become the backup champion. "
            "Tie the tool to a calendar event (prep packet auto-sent) so it becomes habit, not a choice."
        ),
        (
            "Low adoption from analysts",
            "ORGANISATIONAL",
            "Analyst uses the tool but doesn't trust it. Redoes the prep manually anyway. Tool adds a step instead of removing one.",
            "Frame their role as reviewer, not recipient. Make corrections visibly theirs. "
            "Track adoption, not just deployment. If correction rate is low = they trust it. If they're still doing manual prep = investigate why."
        ),
        (
            "Scope creep pressure",
            "ORGANISATIONAL",
            "'Can it also do X?' requests accumulate before Phase 1 is proven. Build gets bloated, trust case gets muddied.",
            "Phase-gate firmly. Every new request goes into a backlog for Phase 2. "
            "The answer in Phase 1 is always: let's prove this works first, then expand."
        ),
    ]

    for i, (title, tag, risk, mitigation) in enumerate(risks):
        top = 1.28 + i * 1.5
        tag_color = RED if tag == "TECHNICAL" else DRED
        # Tag pill
        tag_box = sl.shapes.add_shape(1, Inches(0.2), Inches(top + 0.08), Inches(1.5), Inches(0.28))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = tag_color
        tag_box.line.fill.background()
        tb_tag = txb(sl, 0.22, top + 0.08, 1.45, 0.28)
        para(tb_tag.text_frame, tag, 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        tb1 = txb(sl, 1.85, top, 10.8, 0.38)
        para(tb1.text_frame, title, 13, bold=True, color=DARK)
        tb2 = txb(sl, 0.2, top + 0.42, 5.9, 0.95)
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        para(tf2, "Risk: " + risk, 11, color=GRAY)
        tb3 = txb(sl, 6.3, top + 0.42, 6.8, 0.95)
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        para(tf3, "Mitigation: " + mitigation, 11, color=DARK)


def slide_quickwins_enablement(prs):
    """15 - Quick wins and enablement."""
    sl = blank(prs)
    slide_header(sl, "Quick Wins + Enablement", "First 30 days and what training actually looks like")

    # Left: quick wins
    box(sl, 0.2, 1.28, 6.3, 5.9)
    tb_h = txb(sl, 0.4, 1.35, 5.9, 0.4)
    para(tb_h.text_frame, "Quick wins — first 30 days", 14, bold=True, color=RED)

    wins = [
        "Week 1: First real packet generated for an upcoming meeting.",
        "         Manager reviews it live. One factual thing it got right",
        "         that manual prep would have taken 30 minutes to find.",
        "",
        "Week 2: First open issue caught that wasn't on anyone's radar.",
        "         The Galil Dairy price-cap violation is the template:",
        "         a supplier quote in email that violates contract Section 4.2.",
        "         One concrete catch = trust established faster than any demo.",
        "",
        "Week 3: Correction rate below 20%. Analyst marks fewer fields wrong.",
        "         Show them the trend. That number going down is the proof.",
        "",
        "Week 4: Manager walks into a meeting and says they didn't need",
        "         to do any manual prep. That is the success signal.",
    ]
    bullets(sl, 0.4, 1.85, 5.9, 5.0, wins, size=12)

    # Right: enablement
    box(sl, 6.75, 1.28, 6.3, 5.9)
    tb_h2 = txb(sl, 6.95, 1.35, 5.9, 0.4)
    para(tb_h2.text_frame, "Enablement — what training actually looks like", 14, bold=True, color=RED)

    enablement = [
        "Not a training deck. Not a webinar.",
        "",
        "Session 1 (30 min, Week 1):",
        "  Sit with the analyst. Run the tool together on a real",
        "  upcoming meeting. Walk through every section of the packet.",
        "  Show them the 'Show source' expanders.",
        "  Show them how to flag a correction.",
        "  That is the entire training.",
        "",
        "Session 2 (15 min, Week 3):",
        "  Review the correction log together.",
        "  Which fields did they flag? Any pattern?",
        "  This is the prompt tuning input, framed as their feedback.",
        "",
        "Ongoing:",
        "  One-page reference card: what each section means,",
        "  when to override, who to contact if something looks wrong.",
        "  Non-technical. No jargon.",
    ]
    bullets(sl, 6.95, 1.85, 5.9, 5.0, enablement, size=12)


# ── Build ──────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    slide_title(prs)           # 1
    slide_problem(prs)         # 2
    slide_stakeholders(prs)    # 3  NEW
    slide_week1(prs)           # 4  NEW
    slide_what_built(prs)      # 5
    slide_rejected(prs)        # 6  NEW
    slide_architecture(prs)    # 7
    slide_demo_home(prs)       # 8
    slide_demo_supplier(prs)   # 9
    slide_demo_packet(prs)     # 10
    slide_human_loop(prs)      # 11
    slide_kpis(prs)            # 12
    slide_production(prs)      # 13 NEW
    slide_rollout(prs)         # 14
    slide_risks(prs)           # 14 NEW
    slide_change_management(prs)      # 15
    slide_quickwins_enablement(prs)   # 16 NEW
    slide_roi(prs)             # 17
    slide_whats_next(prs)      # 18

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved {len(prs.slides)} slides -> {OUT}")


if __name__ == "__main__":
    build()
