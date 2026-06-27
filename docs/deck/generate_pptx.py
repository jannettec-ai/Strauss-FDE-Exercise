"""
generate_pptx.py
Generates Strauss FDE Deck.pptx from content extracted from the HTML deck.
Run from project root: python docs/deck/generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
import os
import copy
from lxml import etree

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    """Add a colored rectangle shape."""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = rgb(border_color)
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name='Libre Franklin', font_size=18, bold=False, italic=False,
                color='#221D17', bg_color=None, border_color=None,
                align=PP_ALIGN.LEFT, wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = rgb(bg_color)
    else:
        txBox.fill.background()
    txBox.line.fill.background()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return txBox


def add_textbox_multi(slide, lines, left, top, width, height,
                      font_name='Libre Franklin', font_size=14,
                      color='#221D17', bg_color=None,
                      align=PP_ALIGN.LEFT, wrap=True):
    """Add a textbox with multiple paragraph lines."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = rgb(bg_color)
    else:
        txBox.fill.background()
    txBox.line.fill.background()

    first = True
    for line_data in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align

        if isinstance(line_data, str):
            run = p.add_run()
            run.text = line_data
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = rgb(color)
        elif isinstance(line_data, dict):
            run = p.add_run()
            run.text = line_data.get('text', '')
            run.font.name = line_data.get('font', font_name)
            run.font.size = Pt(line_data.get('size', font_size))
            run.font.bold = line_data.get('bold', False)
            run.font.italic = line_data.get('italic', False)
            run.font.color.rgb = rgb(line_data.get('color', color))
        elif isinstance(line_data, list):
            # Multiple runs in same paragraph
            for seg in line_data:
                run = p.add_run()
                run.text = seg.get('text', '')
                run.font.name = seg.get('font', font_name)
                run.font.size = Pt(seg.get('size', font_size))
                run.font.bold = seg.get('bold', False)
                run.font.italic = seg.get('italic', False)
                run.font.color.rgb = rgb(seg.get('color', color))
    return txBox


def footer(slide, slide_num, total=20, dark=False):
    color = '#7A7165' if dark else '#A89D8E'
    add_textbox(slide, 'STRAUSS PROCUREMENT · FDE EXERCISE · CONFIDENTIAL',
                0.5, 7.1, 10, 0.3, font_name='IBM Plex Mono', font_size=8, color=color)
    add_textbox(slide, f'{slide_num:02d} / {total:02d}',
                12.0, 7.1, 1.1, 0.3, font_name='IBM Plex Mono', font_size=8,
                color=color, align=PP_ALIGN.RIGHT)


def slide_header(slide, section_label, title, dark=False):
    accent = '#CB5238' if dark else '#B23B26'
    title_color = '#F6F1E8' if dark else '#221D17'
    add_textbox(slide, section_label, 0.5, 0.45, 12, 0.3,
                font_name='IBM Plex Mono', font_size=9, color=accent)
    add_textbox(slide, title, 0.5, 0.82, 12.5, 1.0,
                font_name='Newsreader', font_size=32, bold=False, color=title_color,
                wrap=True)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_table(slide, headers, rows, left, top, width, height,
              header_bg='#221D17', header_color='#94897A',
              row_colors=None, font_size=13, header_font_size=11):
    """Add a simple table."""
    num_cols = len(headers)
    num_rows = len(rows) + 1
    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(left), Inches(top), Inches(width), Inches(height)
    ).table

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(header_bg)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.runs[0]
        run.font.name = 'IBM Plex Mono'
        run.font.size = Pt(header_font_size)
        run.font.color.rgb = rgb(header_color)
        run.font.bold = False

    # Data rows
    for ri, row in enumerate(rows):
        bg = row_colors[ri] if row_colors and ri < len(row_colors) else None
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(cell_text)
            if bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(bg)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.runs[0]
            run.font.name = 'Libre Franklin'
            run.font.size = Pt(font_size)
            run.font.color.rgb = rgb('#4A443B')

    return table


# ---------------------------------------------------------------------------
# Presentation setup
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank


# ===========================================================================
# SLIDE 1 — Cover
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#1F1A14')

# Top bar
add_textbox(slide, 'STRAUSS GROUP', 0.5, 0.45, 5, 0.35,
            font_name='IBM Plex Mono', font_size=11, color='#C99A82')
add_textbox(slide, 'FDE HOME EXERCISE / JUNE 2026', 9.5, 0.45, 3.3, 0.35,
            font_name='IBM Plex Mono', font_size=9, color='#8A8072', align=PP_ALIGN.RIGHT)

# Tag
add_textbox(slide, 'Procurement · Forward-Deployed Engineering', 0.5, 2.0, 10, 0.35,
            font_name='IBM Plex Mono', font_size=10, color='#CB5238')

# Main title
add_textbox(slide, 'Negotiation Prep\nIntelligence', 0.5, 2.4, 12.8, 2.2,
            font_name='Newsreader', font_size=72, color='#F7F2E9', wrap=True)

# Subtitle
add_textbox(slide,
            'An AI briefing tool that turns hours of manual meeting prep into a one-page packet '
            '— every fact traceable to its source.',
            0.5, 4.65, 10.5, 0.9,
            font_name='Libre Franklin', font_size=16, color='#B6AC9C', wrap=True)

# Red quote bar
add_rect(slide, 0.5, 5.7, 0.05, 1.1, fill_color='#CB5238')

# Quote
add_textbox(slide,
            '\u201cI want my team spending their time on judgment and relationships. Not on logistics.\u201d',
            0.7, 5.65, 10.5, 0.8,
            font_name='Newsreader', font_size=18, italic=True, color='#E8E1D3', wrap=True)

# Attribution
add_textbox(slide, '— STRAUSS PROCUREMENT MANAGER', 0.7, 6.55, 8, 0.3,
            font_name='IBM Plex Mono', font_size=9, color='#988E7F')

# Bottom right name
add_textbox(slide, 'JANNETTE', 11.5, 6.55, 1.7, 0.3,
            font_name='IBM Plex Mono', font_size=9, color='#8A8072', align=PP_ALIGN.RIGHT)

set_notes(slide,
    "Strauss FDE home exercise — I am presenting a prototype I built in three days. "
    "The brief was: procurement spends too much time on manual meeting prep. "
    "I built one tool that solves exactly that problem. "
    "The title 'Negotiation Prep Intelligence' is deliberate — it is prep intelligence, not a negotiation advisor. "
    "The manager quote on the cover is verbatim from the brief — I am solving the problem they named, not one I invented.")


# ===========================================================================
# SLIDE 2 — The Problem
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '01 — THE PROBLEM',
             'What the procurement team is actually spending time on')

# Stat cards
card_data = [
    ('1.5\u20133 hrs', 'of manual prep per meeting'),
    ('\u20187', 'supplier meetings per month'),
    ('10\u201321 hrs', 'of information work per month, per team'),
]
card_left = 0.5
card_width = 4.1
card_gap = 0.15
for i, (stat, desc) in enumerate(card_data):
    cx = card_left + i * (card_width + card_gap)
    add_rect(slide, cx, 1.95, card_width, 1.1, fill_color='#FBF8F2', border_color='#DDD4C6')
    add_textbox(slide, stat, cx + 0.2, 1.98, card_width - 0.3, 0.65,
                font_name='Newsreader', font_size=38, color='#B23B26')
    add_textbox(slide, desc, cx + 0.2, 2.68, card_width - 0.3, 0.3,
                font_name='Libre Franklin', font_size=14, color='#5C5349', wrap=True)

# Italic quote
add_textbox(slide,
            'The Procurement Manager called it \u201cinformation logistics\u201d \u2014 work that stands '
            'between the team and the judgment call they\u2019re paid to make.',
            0.5, 3.2, 12.5, 0.65,
            font_name='Newsreader', font_size=17, italic=True, color='#3D372F', wrap=True)

# Four failure modes label
add_textbox(slide, 'FOUR WAYS THE MANUAL PROCESS FAILS CONSISTENTLY',
            0.5, 3.9, 12, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#94897A')

# Failure mode grid
failures = [
    ('01', 'Wrong contract version pulled',
     'Missed penalty clause before the negotiation even starts.'),
    ('02', 'Open dispute threads go cold',
     'Unresolved delivery or quality issues forgotten by the next meeting.'),
    ('03', 'Benchmark data goes stale',
     'Manager walks into a cocoa or sugar negotiation blind to where the market moved.'),
    ('04', 'Expired contract not flagged',
     'Team operates on informal terms for months before anyone notices.'),
]
col_w = 6.15
for i, (num, title, desc) in enumerate(failures):
    col = i % 2
    row = i // 2
    fx = 0.5 + col * col_w
    fy = 4.2 + row * 1.0
    add_textbox(slide, num, fx, fy, 0.5, 0.55,
                font_name='Newsreader', font_size=24, color='#B23B26')
    add_textbox(slide, title, fx + 0.55, fy, col_w - 0.7, 0.28,
                font_name='Libre Franklin', font_size=14, bold=True, color='#221D17')
    add_textbox(slide, desc, fx + 0.55, fy + 0.32, col_w - 0.7, 0.55,
                font_name='Libre Franklin', font_size=13, color='#5C5349', wrap=True)

footer(slide, 1)
set_notes(slide,
    "The team runs approximately 7 supplier meetings per month — 14 over the 60-day exercise window — "
    "each needing 1.5 to 3 hours of manual prep. That is up to 21 hours per month of pure information logistics. "
    "I open with the manager's own framing — 'information logistics' — not my framing. "
    "The four failure modes are specific and concrete: wrong contract version, cold dispute threads, "
    "stale benchmark, expired contract. If asked 'how do you know these are the real pain points?' "
    "— they came from watching the process, not describing it.")


# ===========================================================================
# SLIDE 3 — Stakeholder Map
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '02 — STAKEHOLDER MAP',
             'Who uses this, who might resist, and my read on each')

# Table headers
cols = ['Stakeholder', 'What they need', 'My read']
col_widths = [2.6, 3.0, 6.8]  # proportional
table_top = 1.85
table_left = 0.5
table_width = 12.5
table_height = 5.0

tbl = slide.shapes.add_table(6, 3, Inches(table_left), Inches(table_top),
                              Inches(table_width), Inches(table_height)).table

# Set column widths
tbl.columns[0].width = Inches(2.6)
tbl.columns[1].width = Inches(3.0)
tbl.columns[2].width = Inches(6.9)

stakeholders = [
    (
        'Procurement Manager\nPrimary user · Internal champion',
        'Speed + accuracy\nVisible conflict flags',
        'Strongest ally — named the problem themselves. Will adopt fast if the first packet catches something real. Risk: one factual error in the demo and trust is hard to rebuild.'
    ),
    (
        'Procurement Analyst\nDay-to-day user · QA reviewer',
        'Less manual work\nEnough trust to stake prep on it',
        "Highest adoption risk. Won't say they feel threatened — they'll say 'the output is wrong' or 'it adds a review step.' Counter: involve them in Phase 1 QA; make corrections visibly theirs."
    ),
    (
        'IT / Systems\nEnabler for production path',
        'Security review\nData governance\nAPI key management',
        'Not a blocker for the prototype. Critical gating stakeholder for Phase 2 — needs early visibility, not a last-minute ask.'
    ),
    (
        'Legal / Compliance\nRisk owner',
        'AI output never treated as authoritative without human review',
        "Potentially cautious. Frame as 'surfaces for review', not 'interprets contracts.' The correction mechanism is the direct answer to their concern."
    ),
    (
        'Finance / CFO\nCost approver',
        'Measurable ROI in 90 days',
        "Not a day-one stakeholder — bring the ROI case at month 3. Lead with labor-cost avoidance (conservative, auditable), not speculative negotiation savings."
    ),
]

# Header row
header_labels = ['STAKEHOLDER', 'WHAT THEY NEED', 'MY READ']
for ci, hdr in enumerate(header_labels):
    cell = tbl.cell(0, ci)
    cell.text = hdr
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb('#221D17')
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = 'IBM Plex Mono'
    run.font.size = Pt(9)
    run.font.color.rgb = rgb('#94897A')

# Data rows
for ri, (name, need, read) in enumerate(stakeholders):
    row_bg = '#F4EFE7' if ri % 2 == 0 else '#F8F4EE'
    for ci, text in enumerate([name, need, read]):
        cell = tbl.cell(ri + 1, ci)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(row_bg)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = 'Libre Franklin'
        run.font.size = Pt(11)
        run.font.color.rgb = rgb('#4A443B')

footer(slide, 2)
set_notes(slide,
    "Who uses this, who might resist, and my read on each. "
    "The manager is the strongest ally — they named the problem. "
    "The analyst is the highest adoption risk: they will not say they feel threatened, they will say the output is wrong. "
    "I counter by involving them as QA reviewer from day one — their corrections are the feedback loop. "
    "IT is not a blocker for the prototype but is a gating stakeholder for Phase 2. "
    "Legal's concern is addressed by source attribution on every fact. "
    "Finance does not come in until month 3 with actuals.")


# ===========================================================================
# SLIDE 4 — Week 1 Plan
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '03 — WEEK 1 PLAN',
             'Five things to do before touching any technology')

week1 = [
    ('1', 'Shadow a prep session in real time',
     'Sit next to an analyst preparing for a real upcoming meeting. Don\u2019t ask them to describe the process \u2014 watch them do it. The stated process and the actual process are almost always different.'),
    ('2', 'Interview the Procurement Manager separately',
     'Ask: what has gone wrong in a meeting because prep was incomplete? Which suppliers worry you most walking in? What would a perfect prep page look like?'),
    ('3', 'Audit the actual data sources \u2014 not descriptions of them',
     'Request read-only access to one supplier email thread, one contract, and the price benchmark sheet. You cannot design a robust extraction layer without seeing the raw mess.'),
    ('4', 'Map a good packet by working backwards',
     'Pick the three highest-stakes upcoming meetings. Ask the manager: if you had one page of perfect context walking in, what would be on it? Build to that answer, not a generic template.'),
    ('5', 'Identify the first real meeting to demo against',
     'Before building anything, agree on a specific upcoming meeting as the Phase 1 success test. The first demo must be something the manager actually cares about walking into.'),
]

row_h = 0.98
for i, (num, title, desc) in enumerate(week1):
    ry = 1.95 + i * row_h
    add_textbox(slide, num, 0.5, ry, 0.55, 0.7,
                font_name='Newsreader', font_size=32, color='#B23B26')
    add_textbox(slide, title, 1.1, ry, 3.8, 0.32,
                font_name='Libre Franklin', font_size=14, bold=True, color='#221D17', wrap=True)
    add_textbox(slide, desc, 5.1, ry, 8.0, 0.8,
                font_name='Libre Franklin', font_size=12, color='#544C42', wrap=True)

footer(slide, 3)
set_notes(slide,
    "Five things before touching any technology. "
    "The theme: watch the real process, not the described one. "
    "Step one is shadowing a real prep session — the stated process and the actual process are almost always different. "
    "Step three is auditing real data sources, not descriptions of them — "
    "you cannot design a robust extraction layer without seeing the raw mess. "
    "Step five is agreeing a specific meeting to demo against before building anything, "
    "so the first demo is something the manager actually cares about walking into.")


# ===========================================================================
# SLIDE 5 — What I Built
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '04 — WHAT I BUILT', 'One tool. One job. One page per meeting.')

# Body paragraph
add_textbox(slide,
            'An AI-generated one-page briefing, produced automatically before every supplier meeting \u2014 '
            'pulling email history, contract terms and the commodity price benchmark into a single view, '
            'with every claim traceable to its source.',
            0.5, 1.95, 12.5, 1.0,
            font_name='Newsreader', font_size=22, color='#2A241D', wrap=True)

# URL badge
add_rect(slide, 0.5, 3.1, 5.0, 0.42, fill_color='#221D17')
add_textbox(slide, '\u25b6  strauss-procurement.streamlit.app', 0.6, 3.15, 4.8, 0.32,
            font_name='IBM Plex Mono', font_size=13, color='#F1ECE2')

# Three cards
cards = [
    ('EMAILS', 'Summarises six months of supplier threads. Flags open disputes, unanswered threads and price quotes.'),
    ('CONTRACTS', 'Extracts payment terms, volume commitments, penalty clauses and renewal dates from the active contract version.'),
    ('PRICES', 'Adds the live commodity benchmark \u2014 cocoa, sugar, dairy \u2014 so the manager knows where the market sits before walking in.'),
]
card_w = 4.1
cx = 0.5
for label, body in cards:
    add_rect(slide, cx, 3.7, card_w, 2.5, fill_color='#FBF8F2', border_color='#DDD4C6')
    add_textbox(slide, label, cx + 0.2, 3.82, card_w - 0.3, 0.28,
                font_name='IBM Plex Mono', font_size=11, color='#B23B26')
    add_textbox(slide, body, cx + 0.2, 4.15, card_w - 0.35, 1.8,
                font_name='Libre Franklin', font_size=14, color='#403A30', wrap=True)
    cx += card_w + 0.14

footer(slide, 4)
set_notes(slide,
    "One tool, one job, one page per meeting. "
    "An AI-generated briefing pulling email history, contract terms and the commodity benchmark into a single view. "
    "The most important design decision is the last phrase: every claim traceable to its source. "
    "If asked why I built this instead of a pricing dashboard — that was the obvious build and solves an adjacent problem. "
    "The manager's stated pain is information logistics, not pricing accuracy. "
    "The tool is live on Streamlit.")


# ===========================================================================
# SLIDE 6 — Considered & Rejected
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '05 — WHAT I CONSIDERED & REJECTED',
             'Every design decision has a tradeoff \u2014 these were mine')

# Column headers
add_textbox(slide, 'WHAT I CONSIDERED', 0.5, 1.9, 3.2, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#94897A')
add_textbox(slide, 'REJECTED BECAUSE', 3.75, 1.9, 4.8, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#94897A')
add_textbox(slide, 'TRADEOFF ACCEPTED', 8.6, 1.9, 4.7, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#5E7350')

# Divider
add_rect(slide, 0.5, 2.18, 12.5, 0.02, fill_color='#221D17')

rejected = [
    ('Commodity pricing dashboard',
     'Answers an adjacent problem, not the stated pain point. The most obvious build \u2014 every candidate goes here.',
     'Pricing folded in as one input to the packet, not the centrepiece.'),
    ('General \u2018what was agreed\u2019 tracking layer',
     'A system-of-record problem, not an AI problem. Better solved by a database than a model.',
     'Solved indirectly: the packet pulls from accurate, current contract terms.'),
    ('Live MCP connector for pricing in the prototype',
     'A live demo depending on a third-party server responding correctly in the room. Unnecessary risk for zero added insight.',
     'One-time static CSV pull. MCP documented as the production path, not built into the prototype.'),
    ('Vector store + live connectors',
     'Premature at this data volume: 8 suppliers, ~111 emails. Real value arrives at hundreds of suppliers, thousands of emails.',
     'Local extraction against flat files. Vector store called out explicitly as the production-scale answer.'),
    ('Full contract lifecycle management',
     'Too broad a problem for the time available \u2014 and not the time-sink the brief describes.',
     'Contract terms are read and surfaced, not managed end to end.'),
]

row_h = 0.92
for i, (considered, why, tradeoff) in enumerate(rejected):
    ry = 2.22 + i * row_h
    if i > 0:
        add_rect(slide, 0.5, ry, 12.5, 0.01, fill_color='#DBD2C4')
    add_textbox(slide, considered, 0.5, ry + 0.05, 3.1, 0.75,
                font_name='Newsreader', font_size=15, bold=True, color='#221D17', wrap=True)
    add_textbox(slide, why, 3.75, ry + 0.05, 4.7, 0.8,
                font_name='Libre Franklin', font_size=12, color='#4A443B', wrap=True)
    add_rect(slide, 8.55, ry + 0.05, 0.03, 0.7, fill_color='#C3D0B8')
    add_textbox(slide, tradeoff, 8.65, ry + 0.05, 4.35, 0.8,
                font_name='Libre Franklin', font_size=12, color='#3F4A38', wrap=True)

footer(slide, 5)
set_notes(slide,
    "Every design decision has a tradeoff — these were mine. "
    "The common thread: every cut traces back to either the specific pain point in the brief, "
    "the three-day time constraint, or live-demo reliability — not a general 'didn't think of it' gap. "
    "If asked about the commodity dashboard: it solves an adjacent problem, not the stated one. "
    "If asked why no vector store: unjustified at 8 suppliers and 111 emails — right answer at hundreds of suppliers. "
    "If asked why no live MCP connector in the prototype: a live demo depending on a third-party server "
    "in the room is unnecessary risk for zero added insight.")


# ===========================================================================
# SLIDE 7 — How It Works
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '06 — HOW IT WORKS',
             'Prototype vs. production \u2014 a deliberate distinction')

# Left card (prototype)
add_rect(slide, 0.5, 1.85, 6.0, 5.3, fill_color='#FBF8F2', border_color='#DDD4C6')
add_textbox(slide, 'PROTOTYPE', 0.7, 1.98, 3.5, 0.28,
            font_name='IBM Plex Mono', font_size=12, bold=True, color='#221D17')
add_textbox(slide, "what you're seeing today", 2.35, 2.0, 3.5, 0.25,
            font_name='Libre Franklin', font_size=11, italic=True, color='#998E7E')

proto_lines = [
    'SOURCES',
    'Emails \u2192 JSON files (synthetic, ~111) \u00b7 Contracts \u2192 Markdown (10) \u00b7 Prices \u2192 static CSV (yfinance, one-time)',
    '\u2193',
    'extraction.py  \u2014 One Claude API call per supplier \u2014 extracts open issues, price quotes, contract terms, commodity commentary',
    '\u2193',
    'packet_generator.py  \u2014 Assembles packet per meeting ID \u00b7 computes KPIs locally (no API)',
    '\u2193',
    'Streamlit UI  \u2014 Home \u00b7 Suppliers \u00b7 Meeting Prep',
]
proto_colors = ['#B23B26', '#4A443B', '#B6AB99', '#4A443B', '#B6AB99', '#4A443B', '#B6AB99', '#4A443B']
proto_fonts = ['IBM Plex Mono', 'Libre Franklin', 'IBM Plex Mono', 'Libre Franklin',
               'IBM Plex Mono', 'Libre Franklin', 'IBM Plex Mono', 'Libre Franklin']
proto_sizes = [11, 12, 16, 12, 16, 12, 16, 12]

py = 2.38
for line, col, fn, fs in zip(proto_lines, proto_colors, proto_fonts, proto_sizes):
    add_textbox(slide, line, 0.7, py, 5.65, 0.42,
                font_name=fn, font_size=fs, color=col, wrap=True)
    py += 0.56

# Right card (production)
add_rect(slide, 6.8, 1.85, 6.05, 5.3, fill_color='#221D17')
add_textbox(slide, 'PRODUCTION', 7.0, 1.98, 3.5, 0.28,
            font_name='IBM Plex Mono', font_size=12, bold=True, color='#F2ECE0')
add_textbox(slide, 'Phase 2 architecture', 8.85, 2.0, 3.5, 0.25,
            font_name='Libre Franklin', font_size=11, italic=True, color='#A99E8C')

prod_lines = [
    'SOURCES',
    'Emails \u2192 Microsoft Graph (Outlook) \u00b7 Contracts \u2192 SharePoint + OCR \u00b7 Prices \u2192 Bloomberg / ICE live feed',
    '\u2193',
    'Connector layer  \u2014 MCP per source \u2014 swappable backend',
    '\u2193',
    'Storage split  \u2014 Structured store (BigQuery): exact facts \u00b7 Vector store: semantic email search',
    '\u2193',
    'Packet Generator  \u2014 Same logic, real-time data \u00b7 triggered by calendar event \u2192 auto-delivered',
]
prod_colors = ['#CB5238', '#CBC2B2', '#6E6557', '#CBC2B2', '#6E6557', '#CBC2B2', '#6E6557', '#CBC2B2']
prod_fonts = ['IBM Plex Mono', 'Libre Franklin', 'IBM Plex Mono', 'Libre Franklin',
              'IBM Plex Mono', 'Libre Franklin', 'IBM Plex Mono', 'Libre Franklin']
prod_sizes = [11, 12, 16, 12, 16, 12, 16, 12]

py = 2.38
for line, col, fn, fs in zip(prod_lines, prod_colors, prod_fonts, prod_sizes):
    add_textbox(slide, line, 7.0, py, 5.65, 0.42,
                font_name=fn, font_size=fs, color=col, wrap=True)
    py += 0.56

footer(slide, 6)
set_notes(slide,
    "A deliberate distinction between prototype and production — same packet logic throughout, different plumbing. "
    "The prototype runs on flat files: JSON emails, markdown contracts, a static CSV from a one-time yfinance pull. "
    "One Claude API call per supplier extracts open issues, price quotes, contract terms and commodity commentary "
    "into a fixed JSON schema. Production swaps in live connectors: Microsoft Graph for emails, "
    "SharePoint plus OCR for contracts, Bloomberg or ICE for pricing. "
    "MCP is the connector standard — swappable backend per source. "
    "The packet generator is identical in both architectures.")


# ===========================================================================
# SLIDE 8 — Demo Home
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '07 — DEMO \u00b7 HOME DASHBOARD', 'Instant situational awareness')

# Sub-label
add_textbox(slide, 'no API calls \u00b7 loads in <1s', 9.5, 0.88, 3.5, 0.3,
            font_name='IBM Plex Mono', font_size=11, color='#5E7350', align=PP_ALIGN.RIGHT)

# Left bullets
bullets = [
    ('Stat tiles', ' \u2014 upcoming meetings, meetings this week, suppliers at risk, contract issues.'),
    ('Full meeting calendar', ' \u2014 supplier, date, agenda note, geography, urgency badge.'),
    ('Supplier health grid', ' \u2014 relationship status at a glance across all 8 suppliers.'),
    ('Latest discussions', ' \u2014 most recent email per supplier with inline preview.'),
]
by = 1.95
for bold_part, rest in bullets:
    add_textbox(slide, '\u2022  ' + bold_part + rest, 0.5, by, 6.1, 0.48,
                font_name='Libre Franklin', font_size=13, color='#4A443B', wrap=True)
    by += 0.52

# Design decision box
add_rect(slide, 0.5, 4.2, 6.1, 1.4, fill_color='#EEE7DA')
add_rect(slide, 0.5, 4.2, 0.05, 1.4, fill_color='#5E7350')
add_textbox(slide, 'DESIGN DECISION', 0.65, 4.27, 5.5, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#5E7350')
add_textbox(slide,
            'The home page makes ZERO API CALLS. All analytics are computed locally from email and '
            'contract files \u2014 so the manager can load the dashboard anytime, including five minutes '
            'before a meeting, with no cost or latency.',
            0.65, 4.55, 5.8, 0.95,
            font_name='Libre Franklin', font_size=12, color='#3F392F', wrap=True)

# Right: screenshot placeholder
add_rect(slide, 6.8, 1.85, 6.05, 4.8, fill_color='#E8E2D8', border_color='#C8BEB2')
add_textbox(slide, '[Home dashboard screenshot\n\u2014 drag in screenshot]',
            7.5, 3.8, 4.6, 1.0,
            font_name='IBM Plex Mono', font_size=12, color='#A89D8E',
            align=PP_ALIGN.CENTER, wrap=True)
add_textbox(slide, '\u25b6 strauss-procurement.streamlit.app \u00b7 live demo',
            6.8, 6.72, 6.05, 0.28,
            font_name='IBM Plex Mono', font_size=10, color='#A0957F')

footer(slide, 7)
set_notes(slide,
    "Instant situational awareness on load. "
    "The home page makes zero API calls — all analytics computed locally from email and contract files "
    "— so it loads in under a second even if the Anthropic API is unavailable. "
    "Stat tiles show upcoming meetings, suppliers at risk, contract issues. "
    "The full meeting calendar shows supplier, date, agenda note, geography and urgency badge. "
    "The supplier health grid gives relationship status at a glance. "
    "If asked about cost: home page has zero API cost. Only the packet generation step calls Claude.")


# ===========================================================================
# SLIDE 9 — Demo Supplier
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '08 — DEMO \u00b7 SUPPLIER PROFILE',
             'Relationship health, supplier by supplier')

supplier_bullets = [
    ('Searchable directory', ' \u2014 filter by name, category or geography.'),
    ('Relationship health score', ' \u2014 composite of response-time trend + communication volume.'),
    ('Email activity chart', ' \u2014 emails per month across the last 12 months.'),
    ('Responsiveness trend', ' \u2014 early-half vs recent-half response time; getting slower is a warning.'),
    ('Retention signal', ' \u2014 volume trend, days since last contact, interpreted in plain English.'),
    ('Latest discussion preview', ' \u2014 most recent thread with inline read.'),
    ('One-click shortcut', ' to generate the meeting prep packet for this supplier.'),
]
by = 1.92
for bold_part, rest in supplier_bullets:
    add_textbox(slide, '\u2022  ' + bold_part + rest, 0.5, by, 6.1, 0.5,
                font_name='Libre Franklin', font_size=13, color='#4A443B', wrap=True)
    by += 0.55

# Right screenshot placeholder
add_rect(slide, 6.8, 1.85, 6.05, 5.0, fill_color='#E8E2D8', border_color='#C8BEB2')
add_textbox(slide, '[Supplier profile screenshot\n\u2014 drag in screenshot]',
            7.5, 4.0, 4.6, 1.0,
            font_name='IBM Plex Mono', font_size=12, color='#A89D8E',
            align=PP_ALIGN.CENTER, wrap=True)
add_textbox(slide, '\u25b6 live demo \u00b7 searchable per-supplier depth',
            6.8, 6.92, 6.05, 0.28,
            font_name='IBM Plex Mono', font_size=10, color='#A0957F')

footer(slide, 8)
set_notes(slide,
    "Per-supplier relationship health, drill-down view. "
    "The health score is a composite of two signals: response-time trend and communication volume. "
    "The responsiveness trend compares early-half versus recent-half of the email window "
    "— a supplier getting slower is a warning signal before they go silent. "
    "The retention signal interprets days since last contact in plain English. "
    "The one-click shortcut to generate a prep packet closes the loop between supplier context and meeting prep. "
    "If asked about the health score: it is advisory, not a contract risk signal — the packet has the contract data.")


# ===========================================================================
# SLIDE 10 — Demo Packet
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '09 — DEMO \u00b7 MEETING PREP PACKET',
             'The core product \u2014 one click, one page, everything the manager needs')

# Dark badge
add_rect(slide, 0.5, 1.88, 6.1, 0.65, fill_color='#221D17')
add_textbox(slide, '\u26a1  Select meeting \u2192 Generate Packet',
            0.65, 1.93, 5.8, 0.28,
            font_name='Libre Franklin', font_size=14, bold=True, color='#F1ECE2')
add_textbox(slide, 'One Claude API call per supplier \u00b7 returns in ~10\u201315s',
            0.65, 2.22, 5.8, 0.25,
            font_name='Libre Franklin', font_size=11, color='#B6AC9C')

# Packet contents label
add_textbox(slide, 'THE PACKET CONTAINS', 0.5, 2.65, 6, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#94897A')

# 2x4 grid of packet items
items = [
    'Heads-up alert \u2014 most critical issue',
    '4 KPI tiles',
    'Open issues \u2014 severity-coded',
    'Email thread summary',
    'Contract at a glance',
    'Pricing + delta calc',
    'Commodity benchmark \u2014 6-month change',
    '\u201cShow source\u201d on every fact',
]
for idx, item in enumerate(items):
    col = idx % 2
    row = idx // 2
    ix = 0.5 + col * 3.05
    iy = 2.95 + row * 0.55
    add_textbox(slide, item, ix, iy, 2.9, 0.45,
                font_name='Libre Franklin', font_size=12, color='#4A443B', wrap=True)

# Right panel: Demo case
add_rect(slide, 6.8, 1.85, 6.05, 5.3, fill_color='#FBF8F2', border_color='#DDD4C6')
add_textbox(slide, 'Demo case \u00b7 Galil Dairy', 7.0, 1.98, 4.0, 0.32,
            font_name='Newsreader', font_size=18, bold=True, color='#221D17')
add_textbox(slide, 'GENERATED PACKET', 10.3, 2.02, 2.3, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#A0957F', align=PP_ALIGN.RIGHT)

# Red alert
add_rect(slide, 6.9, 2.4, 5.85, 1.2, fill_color='#F7ECE9', border_color='#E6C7BF')
add_textbox(slide, '\u2022  Price discrepancy flagged', 7.05, 2.47, 5.5, 0.28,
            font_name='Libre Franklin', font_size=13, bold=True, color='#9A2E1C')
add_textbox(slide,
            'Supplier quoted \u20aa15.50/kg in email; contract \u00a74.2 caps at \u20aa14.20/kg. '
            'No written amendment exists.  Delta +\u20aa1.30/kg (+9.2%)',
            7.05, 2.77, 5.5, 0.72,
            font_name='Libre Franklin', font_size=12, color='#4A443B', wrap=True)

# Amber alert
add_rect(slide, 6.9, 3.68, 5.85, 1.1, fill_color='#FAF2E2', border_color='#EAD7AE')
add_textbox(slide, '\u2022  Unanswered thread flagged', 7.05, 3.75, 5.5, 0.28,
            font_name='Libre Franklin', font_size=13, bold=True, color='#8A6212')
add_textbox(slide,
            'Yael Shapira requested alignment; David Cohen proposed a call \u2014 no documented outcome.',
            7.05, 4.05, 5.5, 0.6,
            font_name='Libre Franklin', font_size=12, color='#4A443B', wrap=True)

# Source note
add_textbox(slide,
            '\U0001f4ce  Every fact carries a \u201cShow source\u201d expander pointing to the exact email or contract section.',
            7.0, 4.88, 5.7, 0.65,
            font_name='Libre Franklin', font_size=12, color='#544C42', wrap=True)

footer(slide, 9)
set_notes(slide,
    "The core product — one click, one Claude API call, returns in 10 to 15 seconds. "
    "The Galil Dairy case is the money shot for the demo: supplier quoted 15.50 shekel per kilo in email, "
    "contract section 4.2 caps at 14.20, and there is no written amendment on file. "
    "The packet surfaced this discrepancy with source attribution — the exact email date and contract section. "
    "Every field in the packet has a 'Show source' expander. "
    "The 24-hour cache means the second load is instant — no repeated API cost. "
    "If asked about accuracy: source attribution is the answer — every claim is one click from its evidence.")


# ===========================================================================
# SLIDE 11 — Human in the Loop
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '10 — THE HUMAN IN THE LOOP',
             'What stops this from being trusted blindly')

hitl_cards = [
    ('01', 'Source attribution on every field',
     'Every number, date and clause cites exactly where it came from \u2014 email date + sender, '
     'or contract section. One click to verify. Nothing is presented as settled ground truth.',
     False),
    ('02', 'Explicit conflict flagging',
     'When email and contract disagree, the packet surfaces both and flags the conflict \u2014 '
     'it never silently picks one. The Galil price discrepancy and Siam Sugar unit ambiguity '
     'are live examples in the demo.',
     False),
    ('03', 'Correction tracking built in',
     'The manager can flag any field as wrong directly in the packet view. '
     'Corrections save to a log correction_log.csv \u2014 the extraction-quality feedback loop.',
     False),
    ('04', 'The tool never acts',
     'No emails sent, no contracts updated, no decisions taken. Every action stays with the human. '
     'The tool removes the assembly work in front of the judgment call \u2014 '
     'it does not replace the judgment call.',
     True),
]

positions = [
    (0.5, 1.85, 6.15, 2.55),
    (6.7, 1.85, 6.15, 2.55),
    (0.5, 4.5, 6.15, 2.55),
    (6.7, 4.5, 6.15, 2.55),
]

for (num, title, body, dark), (lx, ty, w, h) in zip(hitl_cards, positions):
    bg = '#221D17' if dark else '#FBF8F2'
    border = '#221D17' if dark else '#DDD4C6'
    title_col = '#F5F0E6' if dark else '#221D17'
    num_col = '#CB5238' if dark else '#B23B26'
    body_col = '#C7BEAE' if dark else '#4A443B'

    add_rect(slide, lx, ty, w, h, fill_color=bg, border_color=border)
    add_textbox(slide, num, lx + 0.2, ty + 0.18, 0.5, 0.3,
                font_name='IBM Plex Mono', font_size=13, color=num_col)
    add_textbox(slide, title, lx + 0.75, ty + 0.15, w - 0.9, 0.35,
                font_name='Newsreader', font_size=18, bold=True, color=title_col, wrap=True)
    add_textbox(slide, body, lx + 0.2, ty + 0.6, w - 0.35, h - 0.75,
                font_name='Libre Franklin', font_size=13, color=body_col, wrap=True)

footer(slide, 10)
set_notes(slide,
    "What stops this from being trusted blindly — four concrete mechanisms, not a general disclaimer. "
    "Source attribution means nothing is presented as settled ground truth. "
    "Conflict flagging means when email and contract disagree, both surfaces — the tool never silently picks one. "
    "The Galil price discrepancy and Siam Sugar unit ambiguity are live examples in the demo. "
    "The correction mechanism is built into the UI as a first-class feature, not an afterthought "
    "— corrections feed the trust signal. "
    "The tool never acts: no emails sent, no contracts updated, every decision stays with the human.")


# ===========================================================================
# SLIDE 12 — KPIs
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '11 — KPIs', 'Two tiers, both tracked \u2014 adoption, then quality')

# Tier 1 label
add_textbox(slide, 'TIER 1 \u00b7 ADOPTION', 0.5, 1.88, 3.5, 0.25,
            font_name='IBM Plex Mono', font_size=10, bold=True, color='#221D17')
add_textbox(slide, 'is it being used?', 3.85, 1.9, 3.0, 0.22,
            font_name='Libre Franklin', font_size=11, italic=True, color='#998E7E')

tier1 = [
    ('Packets generated',
     'Total prep packets produced. The primary adoption signal \u2014 zero means the tool isn\u2019t being used.',
     'total_packets'),
    ('Unique suppliers queried',
     'Distinct suppliers with a packet generated. Breadth of adoption across the roster.',
     'unique_suppliers'),
    ('Avg generation time',
     'Seconds per packet. Tracks performance \u2014 target <15s; a regression signal if it grows.',
     'avg_duration_sec'),
    ('Daily generation volume',
     'Packets/day over a rolling 30 days. Shows whether use is growing, flat, or fading after curiosity.',
     'daily_breakdown'),
    ('Per-supplier frequency',
     'Which suppliers get prepped most. Low frequency on a high-risk supplier = an adoption gap.',
     'supplier_breakdown'),
]

card_w = 2.45
cx = 0.5
for name, desc, col in tier1:
    add_rect(slide, cx, 2.18, card_w, 1.65, fill_color='#FBF8F2', border_color='#DDD4C6')
    add_textbox(slide, name, cx + 0.15, 2.25, card_w - 0.25, 0.38,
                font_name='Libre Franklin', font_size=12, bold=True, color='#221D17', wrap=True)
    add_textbox(slide, desc, cx + 0.15, 2.65, card_w - 0.25, 0.85,
                font_name='Libre Franklin', font_size=10, color='#5C5349', wrap=True)
    add_textbox(slide, col, cx + 0.15, 3.62, card_w - 0.25, 0.2,
                font_name='IBM Plex Mono', font_size=8, color='#A0957F')
    cx += card_w + 0.12

# Tier 2 label
add_textbox(slide, 'TIER 2 \u00b7 EXTRACTION QUALITY', 0.5, 4.05, 5.0, 0.25,
            font_name='IBM Plex Mono', font_size=10, bold=True, color='#B23B26')
add_textbox(slide, 'is it worth using?', 5.3, 4.07, 3.0, 0.22,
            font_name='Libre Franklin', font_size=11, italic=True, color='#998E7E')

tier2 = [
    ('1', 'Response Time',
     'Avg days from Strauss outbound to the first supplier reply.',
     '<5d healthy \u00b7 >14d = relationship risk', 'response_time_days'),
    ('2', 'Open Issues',
     'Unresolved disputes, threads and contract gaps per supplier.',
     'Target 0 \u2014 each one costs prep time', 'open_issues_count'),
    ('3', 'Price Delta',
     'Latest email quote vs contract base price \u2014 % and absolute.',
     'Alert if >0% without a written amendment', 'price_delta_pct'),
    ('4', 'Days to Renewal',
     'Calendar days to the contract renewal date.',
     'Red <30d \u00b7 amber 30\u201390d', 'days_to_renewal'),
    ('5', 'Correction Rate',
     '% of AI fields the manager flags wrong per packet.',
     'Target <10% by day 90 \u00b7 high = retune prompt', 'correction_rate_pct'),
]

cx = 0.5
for n, name, desc, bench, col in tier2:
    add_rect(slide, cx, 4.35, card_w, 2.35, fill_color='#FBF8F2', border_color='#DDD4C6')
    add_rect(slide, cx, 4.35, card_w, 0.06, fill_color='#B23B26')
    add_textbox(slide, f'\u00a7{n}  {name}', cx + 0.15, 4.45, card_w - 0.25, 0.4,
                font_name='Libre Franklin', font_size=12, bold=True, color='#221D17', wrap=True)
    add_textbox(slide, desc, cx + 0.15, 4.9, card_w - 0.25, 0.85,
                font_name='Libre Franklin', font_size=10, color='#5C5349', wrap=True)
    add_textbox(slide, bench, cx + 0.15, 5.82, card_w - 0.25, 0.4,
                font_name='Libre Franklin', font_size=10, color='#7A6E2E', wrap=True)
    add_textbox(slide, col, cx + 0.15, 6.35, card_w - 0.25, 0.2,
                font_name='IBM Plex Mono', font_size=8, color='#A0957F')
    cx += card_w + 0.12

footer(slide, 11)
set_notes(slide,
    "Two tiers, both tracked. Tier 1 asks: is it being used? "
    "Packets generated, unique suppliers queried, generation time, daily volume, per-supplier frequency. "
    "Tier 2 asks: is it worth using? The five metrics from the metrics dictionary "
    "— response time, open issues, price delta, days to renewal, correction rate. "
    "The correction rate is the one that matters most: it is the trust signal. "
    "If it stays high, the extraction prompt needs tuning. "
    "If it drops below 10 percent, the team trusts the output.")


# ===========================================================================
# SLIDE 13 — KPI Data Inputs
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '12 — KPI DATA INPUTS',
             'What each metric reads, where it comes from, who computes it')

# Table
kpi_headers = ['KPI', 'Input fields', 'Source system (prototype)', 'Who computes it']
kpi_rows = [
    ('\u00a71  Response Time',
     'email.date \u00b7 email.from \u00b7 subject',
     'data/emails/<supplier>.json \u2014 111 emails',
     'Python (local) \u00b7 no API \u00b7 supplier_analytics.py'),
    ('\u00a72  Open Issues',
     'Full email body \u00b7 contract clause text \u00b7 issue type / status',
     'Same email JSONs + data/contracts/*.md, fed to Claude',
     'Claude API \u00b7 extraction.py \u00b7 structured JSON'),
    ('\u00a73  Price Delta',
     'latest_price_quoted.value / .unit \u00b7 contract_base_price_numeric',
     'Price: Claude extraction \u00b7 Base: contract Pricing section',
     'Claude (price) + Python delta calc \u00b7 packet_generator.py'),
    ('\u00a74  Days to Renewal',
     'contract.renewal_date \u00b7 system date \u00b7 contract.status',
     'data/contracts/*.md \u2014 Renewal Date field',
     'Python (local) \u00b7 date arithmetic only \u00b7 no API'),
    ('\u00a75  Correction Rate',
     'Checkbox state per field \u00b7 8 trackable fields / packet',
     'Streamlit session state \u2192 correction_log.csv',
     'UI interaction + Python writes log \u00b7 no model'),
]

tbl = slide.shapes.add_table(6, 4,
                              Inches(0.5), Inches(2.0),
                              Inches(12.5), Inches(4.5)).table

tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(3.2)
tbl.columns[2].width = Inches(3.8)
tbl.columns[3].width = Inches(3.5)

hdr_labels = ['KPI', 'INPUT FIELDS', 'SOURCE SYSTEM (PROTOTYPE)', 'WHO COMPUTES IT']
for ci, hdr in enumerate(hdr_labels):
    cell = tbl.cell(0, ci)
    cell.text = hdr
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb('#221D17')
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = 'IBM Plex Mono'
    run.font.size = Pt(9)
    run.font.color.rgb = rgb('#94897A')

for ri, row in enumerate(kpi_rows):
    bg = '#F4EFE7' if ri % 2 == 0 else '#F8F4EE'
    for ci, text in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(bg)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = 'Libre Franklin'
        run.font.size = Pt(11)
        run.font.color.rgb = rgb('#4A443B')

# Note
add_textbox(slide,
            'NOTE  \u00a71, \u00a73, \u00a74 are pure Python (no API cost). '
            '\u00a72 is the only metric that needs a Claude call. '
            '\u00a75 is populated only when the manager actively saves corrections.',
            0.5, 6.6, 12.5, 0.45,
            font_name='Libre Franklin', font_size=12, color='#544C42', wrap=True)

footer(slide, 12)
set_notes(slide,
    "Exactly what data each metric reads, where it comes from, and who computes it. "
    "The key architectural point: sections 1, 3 and 4 are pure Python with no API cost "
    "— date arithmetic, JSON field reads. Only section 2 (open issues) needs a Claude call. "
    "Correction rate is populated only when the manager actively saves a flag. "
    "If asked who defined these metrics: they are in a metrics dictionary at .claude/rules/metrics.md "
    "— single source of truth, not redefined inline. This matters for auditability.")


# ===========================================================================
# SLIDE 14 — Production Architecture
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '13 — PRODUCTION ARCHITECTURE',
             'How the prototype becomes a production system')

# Three columns
col_labels = ['Data Sources', 'Processing Layer', 'Outputs + Monitoring']
col_x = [0.5, 4.85, 9.2]
col_w = 3.9

for i, label in enumerate(col_labels):
    add_textbox(slide, label.upper(), col_x[i], 1.88, col_w, 0.22,
                font_name='IBM Plex Mono', font_size=9, color='#94897A')

# Arrows
add_textbox(slide, '\u2192', 4.55, 3.6, 0.3, 0.4,
            font_name='IBM Plex Mono', font_size=22, color='#C0B6A4')
add_textbox(slide, '\u2192', 8.9, 3.6, 0.3, 0.4,
            font_name='IBM Plex Mono', font_size=22, color='#C0B6A4')

# Data sources cards
ds_cards = [
    ('Emails', 'Microsoft Graph API \u00b7 Outlook / Exchange \u00b7 shared procurement mailbox \u00b7 push subscription on new mail'),
    ('Contracts', 'SharePoint connector \u00b7 PDF extraction + OCR \u00b7 version detection (active vs draft)'),
    ('Prices', 'Bloomberg / ICE / Refinitiv \u00b7 structured time-series feed \u00b7 normalised to a common schema'),
]
cy = 2.18
for title, body in ds_cards:
    add_rect(slide, 0.5, cy, col_w, 1.5, fill_color='#FBF8F2', border_color='#DDD4C6')
    add_textbox(slide, title, 0.68, cy + 0.1, col_w - 0.25, 0.28,
                font_name='Libre Franklin', font_size=13, bold=True, color='#221D17')
    add_textbox(slide, body, 0.68, cy + 0.42, col_w - 0.25, 0.95,
                font_name='Libre Franklin', font_size=11, color='#544C42', wrap=True)
    cy += 1.6

# Processing layer cards
pl_cards = [
    ('MCP connector layer', 'Standardised tool interface \u00b7 swappable backend per source'),
    ('Storage split', 'Structured (BigQuery): exact facts \u2014 contract terms, KPI values, price points.\nVector (Vertex / Pinecone): chunked emails + clause language, semantic search at scale.'),
    ('Packet Generator', 'Same logic as the prototype \u00b7 triggered by calendar event 24h before each meeting'),
]
cy = 2.18
for title, body in pl_cards:
    add_rect(slide, 4.85, cy, col_w, 1.5, fill_color='#FBF8F2', border_color='#DDD4C6')
    add_textbox(slide, title, 5.03, cy + 0.1, col_w - 0.25, 0.28,
                font_name='Libre Franklin', font_size=13, bold=True, color='#221D17')
    add_textbox(slide, body, 5.03, cy + 0.42, col_w - 0.25, 0.95,
                font_name='Libre Franklin', font_size=11, color='#544C42', wrap=True)
    cy += 1.6

# Output cards
out_cards = [
    ('Procurement Manager UI', 'Streamlit or Teams tab \u00b7 packet auto-delivered \u00b7 source attribution + conflict flags', False),
    ('FDE Metrics Pipeline', 'Generation events Pub/Sub \u2192 BigQuery events table \u00b7 KPI values per packet stored \u00b7 correction log aggregated \u00b7 Looker dashboard', True),
    ('Prompt tuning cycle', 'Monthly: correction_log patterns \u2192 extraction prompt update \u2192 re-run accuracy benchmark', False),
]
cy = 2.18
for title, body, dark in out_cards:
    bg = '#221D17' if dark else '#FBF8F2'
    border = '#221D17' if dark else '#DDD4C6'
    tc = '#F2ECE0' if dark else '#221D17'
    bc = '#CBC2B2' if dark else '#544C42'
    add_rect(slide, 9.2, cy, col_w, 1.5, fill_color=bg, border_color=border)
    add_textbox(slide, title, 9.38, cy + 0.1, col_w - 0.25, 0.28,
                font_name='Libre Franklin', font_size=13, bold=True, color=tc)
    add_textbox(slide, body, 9.38, cy + 0.42, col_w - 0.25, 0.95,
                font_name='Libre Franklin', font_size=11, color=bc, wrap=True)
    cy += 1.6

# Prototype delta note
add_textbox(slide,
            'PROTOTYPE DELTA: Flat JSON files instead of connectors \u00b7 local extraction instead of '
            'a vector store \u00b7 SQLite instead of BigQuery. Every simplification is documented and justified.',
            0.5, 6.88, 12.5, 0.38,
            font_name='IBM Plex Mono', font_size=9, color='#94897A', wrap=True)

footer(slide, 13)
set_notes(slide,
    "How the prototype becomes a production system — three columns, left to right. "
    "Data sources: Microsoft Graph for emails, SharePoint plus OCR for contracts, Bloomberg or ICE for pricing. "
    "Processing layer: MCP connector per source with a swappable backend, "
    "a storage split between BigQuery for exact facts and a vector store for semantic email search. "
    "Outputs: Streamlit or Teams tab, auto-delivered 24 hours before each meeting. "
    "The metrics pipeline feeds a monthly prompt-tuning cycle. "
    "Every prototype simplification is documented: flat files instead of connectors, local extraction instead of a vector store.")


# ===========================================================================
# SLIDE 15 — Rollout Plan
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '14 — ROLLOUT PLAN', 'Three phases \u2014 trust first, scale second')

phases = [
    ('PHASE 1', 'Weeks 1\u20134', 'Controlled pilot', [
        '2\u20133 upcoming meetings with one analyst as primary user.',
        'Analyst runs the tool, reviews output, flags corrections.',
        'Track correction rate weekly \u2014 target <20% by week 4.',
        'One structured debrief after each meeting: what did the packet get right?',
        'Success gate: Packet catches one real issue the analyst would have missed.',
    ]),
    ('PHASE 2', 'Weeks 5\u20138', 'Expanded pilot', [
        'All 8 suppliers in scope; full procurement team using it.',
        'IT review: security, data access, API key governance.',
        'Calendar integration auto-triggers packet generation 24h before each meeting.',
        'Add live sources: Graph API for emails, SharePoint for contracts.',
        'Success gate: Team uses the packet without being reminded; correction rate <10%.',
    ]),
    ('PHASE 3', 'Weeks 9\u201312', 'Handover', [
        'Ownership transferred to the internal champion (Procurement Manager).',
        'ROI case presented to Finance / CFO with 90-day actuals.',
        'Prompt-tuning backlog based on correction_log.csv patterns.',
        'Roadmap: supplier segmentation, auto-renewal alerts, negotiation history.',
        'Outcome: A tool the team owns \u2014 not one they depend on me for.',
    ]),
]

col_w = 4.0
cx = 0.5
for phase_label, weeks, subtitle, items in phases:
    add_textbox(slide, phase_label, cx, 1.9, col_w, 0.3,
                font_name='IBM Plex Mono', font_size=12, bold=True, color='#B23B26')
    add_textbox(slide, weeks, cx, 2.22, col_w, 0.25,
                font_name='IBM Plex Mono', font_size=10, color='#94897A')
    add_textbox(slide, subtitle, cx, 2.52, col_w, 0.28,
                font_name='Newsreader', font_size=16, color='#221D17')

    add_rect(slide, cx, 2.88, col_w, 4.1, fill_color='#FBF8F2', border_color='#DDD4C6')
    iy = 3.0
    for item in items:
        add_textbox(slide, '\u2022  ' + item, cx + 0.15, iy, col_w - 0.25, 0.6,
                    font_name='Libre Franklin', font_size=12, color='#4A443B', wrap=True)
        iy += 0.72

    cx += col_w + 0.28

footer(slide, 14)
set_notes(slide,
    "Three phases — trust first, scale second. "
    "Phase 1 is a controlled pilot with one analyst and two or three suppliers. "
    "The success gate is not adoption rate — it is whether the packet catches one real issue "
    "the analyst would have missed. One concrete catch builds trust faster than any demo. "
    "Phase 2 expands to the full team with auto-triggering from the calendar and live data connectors. "
    "Phase 3 hands ownership to the internal champion with the 90-day ROI case for the CFO. "
    "If asked why not roll out to everyone at once: a high error rate across 14 meetings simultaneously "
    "destroys trust before it is built.")


# ===========================================================================
# SLIDE 16 — Risks & Challenges
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '15 — RISKS & CHALLENGES',
             'What could go wrong \u2014 and how I\u2019d handle it')

risks = [
    ('TECHNICAL', 'Extraction accuracy drift',
     'RISK: Claude returns wrong contract terms or misses a dispute. Manager walks in with bad data.',
     'MITIGATION: Source attribution on every field makes errors visible. The correction log tracks error '
     'rate by field and supplier. Monthly prompt-tuning cycle against correction_log.csv. '
     'Phase 1 KPI: correction rate <10% by week 4.'),
    ('ORGANISATIONAL', 'Losing the internal champion',
     'RISK: Procurement Manager moves roles or loses interest. Tool loses its sponsor before Phase 2.',
     'MITIGATION: Involve the analyst as co-owner from day one \u2014 they become the backup champion. '
     'Tie the tool to a calendar event (packet auto-sent) so it becomes habit, not a choice.'),
    ('ORGANISATIONAL', 'Low adoption from analysts',
     "RISK: Analyst uses the tool but doesn't trust it \u2014 redoes prep manually anyway. "
     'The tool adds a step instead of removing one.',
     'MITIGATION: Frame their role as reviewer, not recipient. Make corrections visibly theirs. '
     'Track adoption, not just deployment: low correction rate = trust; still doing manual prep = investigate why.'),
    ('ORGANISATIONAL', 'Scope creep pressure',
     'RISK: "Can it also do X?" requests pile up before Phase 1 is proven. '
     'The build bloats, the trust case gets muddied.',
     'MITIGATION: Phase-gate firmly. Every new request goes to a Phase 2 backlog. '
     'The Phase 1 answer is always: let\u2019s prove this works first, then expand.'),
]

positions_r = [
    (0.5, 1.85, 6.15, 2.6),
    (6.7, 1.85, 6.15, 2.6),
    (0.5, 4.55, 6.15, 2.6),
    (6.7, 4.55, 6.15, 2.6),
]

for (category, title, risk_text, mitigation), (lx, ty, w, h) in zip(risks, positions_r):
    add_rect(slide, lx, ty, w, h, fill_color='#FBF8F2', border_color='#DDD4C6')
    add_textbox(slide, category, lx + 0.2, ty + 0.12, w - 0.3, 0.22,
                font_name='IBM Plex Mono', font_size=9, color='#94897A')
    add_textbox(slide, title, lx + 0.2, ty + 0.36, w - 0.3, 0.32,
                font_name='Newsreader', font_size=17, bold=True, color='#221D17', wrap=True)
    add_textbox(slide, risk_text, lx + 0.2, ty + 0.72, w - 0.3, 0.65,
                font_name='Libre Franklin', font_size=11, color='#9A2E1C', wrap=True)
    add_textbox(slide, mitigation, lx + 0.2, ty + 1.42, w - 0.3, 1.05,
                font_name='Libre Franklin', font_size=11, color='#4A443B', wrap=True)

footer(slide, 15)
set_notes(slide,
    "Technical and organisational — the honest read. "
    "The technical risk is extraction accuracy drift: Claude returns wrong contract terms "
    "and the manager walks in with bad data. "
    "The source attribution and correction log are the direct mitigations — errors are visible, not silent. "
    "But most of the real risk is organisational: losing the internal champion, low analyst adoption, and scope creep. "
    "The analyst objection is the one that decides whether the tool gets used "
    "— they will not say they feel threatened, they will say the output is wrong. "
    "Involve them in QA from day one.")


# ===========================================================================
# SLIDE 17 — Change Management
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '16 — CHANGE MANAGEMENT',
             'Building it is 20% of the job. Adoption is the other 80%.')

cm_rows = [
    ('\"The output is wrong.\"',
     'PROCUREMENT ANALYST\nHighest adoption risk',
     'Involve them in Phase 1 QA from day one. Make the correction feature visibly theirs. '
     'Frame their role as reviewer, not recipient \u2014 the tool makes them look more prepared, not redundant.'),
    ('\"I don\u2019t trust AI to read a contract.\"',
     'LEGAL / COMPLIANCE',
     'Every fact carries a source citation. The design surfaces facts for human review \u2014 '
     'it does not interpret and act. Show them the correction log: the system has a visible trust mechanism.'),
    ('\"What does this cost, and who owns it?\"',
     'FINANCE / IT',
     'Phase 1 cost is API calls only (~$15/month). Nothing touches live systems; '
     'the static-file architecture minimises the security surface. ROI case ready for month 3.'),
]

# Column headers
add_textbox(slide, 'OBJECTION', 0.5, 1.88, 4.0, 0.22,
            font_name='IBM Plex Mono', font_size=9, color='#94897A')
add_textbox(slide, 'WHO', 4.55, 1.88, 2.5, 0.22,
            font_name='IBM Plex Mono', font_size=9, color='#94897A')
add_textbox(slide, 'HOW I HANDLE IT', 7.1, 1.88, 6.0, 0.22,
            font_name='IBM Plex Mono', font_size=9, color='#94897A')

add_rect(slide, 0.5, 2.12, 12.5, 0.03, fill_color='#221D17')

row_h = 1.55
for i, (objection, who, handle) in enumerate(cm_rows):
    ry = 2.18 + i * row_h
    if i > 0:
        add_rect(slide, 0.5, ry, 12.5, 0.01, fill_color='#DBD2C4')
    add_textbox(slide, objection, 0.5, ry + 0.12, 3.8, 0.65,
                font_name='Newsreader', font_size=18, italic=True, color='#221D17', wrap=True)
    add_textbox(slide, who, 4.55, ry + 0.12, 2.3, 0.65,
                font_name='IBM Plex Mono', font_size=10, bold=True, color='#B23B26', wrap=True)
    add_textbox(slide, handle, 7.1, ry + 0.08, 5.9, 1.3,
                font_name='Libre Franklin', font_size=13, color='#4A443B', wrap=True)

footer(slide, 16)
set_notes(slide,
    "Building it is 20 percent of the job. Adoption is the other 80. "
    "Three objections I expect, who raises each, and exactly how I answer. "
    "The analyst objection — 'the output is wrong' — is the one that decides whether the tool gets used. "
    "I counter by making them the reviewer, not the recipient, from day one. "
    "Their corrections are visibly theirs. "
    "Legal's concern is addressed by source attribution: the tool surfaces for review, it does not interpret and act. "
    "Finance's question is about cost — Phase 1 API cost is approximately $5 per month, "
    "nothing touches live systems.")


# ===========================================================================
# SLIDE 18 — Quick Wins + Enablement
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '17 — QUICK WINS + ENABLEMENT',
             'First 30 days, and what training actually looks like')

# Left half — Quick wins
add_textbox(slide, 'QUICK WINS', 0.5, 1.9, 6.0, 0.25,
            font_name='IBM Plex Mono', font_size=9, bold=True, color='#221D17')

quick_wins = [
    ('WEEK 1',
     'First real packet generated for an upcoming meeting. Manager reviews it live \u2014 one factual thing '
     'it got right that manual prep would have taken 30 minutes to find.'),
    ('WEEK 2',
     'First open issue caught that wasn\u2019t on anyone\u2019s radar. The Galil Dairy price-cap violation '
     'is the template. One concrete catch establishes trust faster than any demo.'),
    ('WEEK 3',
     'Correction rate below 20%. Analyst marks fewer fields wrong. Show them the trend \u2014 '
     'that number going down is the proof.'),
    ('WEEK 4',
     'Manager walks into a meeting and says they didn\u2019t need to do any manual prep. '
     'That is the success signal.'),
]

wy = 2.2
for week, desc in quick_wins:
    add_textbox(slide, week, 0.5, wy, 1.2, 0.25,
                font_name='IBM Plex Mono', font_size=10, bold=True, color='#B23B26')
    add_textbox(slide, desc, 1.75, wy, 4.6, 0.7,
                font_name='Libre Franklin', font_size=12, color='#4A443B', wrap=True)
    wy += 0.88

# Right half — Enablement
add_textbox(slide, 'ENABLEMENT', 6.9, 1.9, 6.0, 0.25,
            font_name='IBM Plex Mono', font_size=9, bold=True, color='#221D17')
add_textbox(slide, 'Not a training deck. Not a webinar.', 6.9, 2.2, 5.9, 0.3,
            font_name='Newsreader', font_size=15, italic=True, color='#544C42')

enablement = [
    ('Session 1 \u00b7 30 min \u00b7 Week 1',
     'Sit with the analyst. Run the tool together on a real upcoming meeting. Walk through every section. '
     'Show the \u201cShow source\u201d expanders and how to flag a correction. That is the entire training.'),
    ('Session 2 \u00b7 15 min \u00b7 Week 3',
     'Review the correction log together. Which fields did they flag? Any pattern? '
     'This is the prompt-tuning input, framed as their feedback.'),
    ('Ongoing',
     'A one-page reference card: what each section means, when to override, who to contact if something looks wrong. '
     'Non-technical. No jargon.'),
]

ey = 2.6
for session, desc in enablement:
    add_textbox(slide, session, 6.9, ey, 5.9, 0.25,
                font_name='IBM Plex Mono', font_size=10, bold=True, color='#5E7350')
    add_textbox(slide, desc, 6.9, ey + 0.28, 5.9, 0.9,
                font_name='Libre Franklin', font_size=12, color='#4A443B', wrap=True)
    ey += 1.22

footer(slide, 17)
set_notes(slide,
    "First 30 days and what training actually looks like. "
    "The enablement is deliberately not a deck or a webinar "
    "— it is two short sessions sitting next to the analyst on a real meeting. "
    "Session 1 is 30 minutes in week 1: run the tool together on a real upcoming meeting, "
    "walk through every section, show the source expanders and how to flag a correction. "
    "That is the entire training. "
    "The Galil catch in week 2 is the trust-builder "
    "— one concrete catch establishes trust faster than any demo. "
    "Week 4 success signal: manager says they did not do any manual prep.")


# ===========================================================================
# SLIDE 19 — The ROI Case
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#F4EFE7')

slide_header(slide, '18 — THE ROI CASE',
             'Lead with what\u2019s measurable, then the multiplier')

# Three column layout
# LEFT — conservative number
add_rect(slide, 0.5, 1.85, 3.95, 5.25, fill_color='#FBF8F2', border_color='#DDD4C6')
add_textbox(slide, 'THE CONSERVATIVE NUMBER', 0.65, 1.95, 3.7, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#94897A')
add_textbox(slide, 'Labor cost avoidance', 0.65, 2.22, 3.7, 0.28,
            font_name='Newsreader', font_size=16, color='#221D17')

roi_left = [
    ('Meetings / month:', '7'),
    ('Prep time / meeting:', '135 \u2192 15 min'),
    ('Time saved:', '14 hrs / mo'),
    ('Analyst cost [confirm w/ HR]:', '\u20aa200 / hr'),
    ('API cost:', '~$5/mo \u00b7 \u22480'),
]
ly = 2.58
for label, val in roi_left:
    add_textbox(slide, label, 0.65, ly, 2.3, 0.28,
                font_name='Libre Franklin', font_size=11, color='#5C5349')
    add_textbox(slide, val, 2.95, ly, 1.3, 0.28,
                font_name='Libre Franklin', font_size=11, bold=True, color='#221D17',
                align=PP_ALIGN.RIGHT)
    ly += 0.38

add_rect(slide, 0.65, ly + 0.05, 3.65, 0.02, fill_color='#DDD4C6')
add_textbox(slide, 'TOTAL', 0.65, ly + 0.12, 2.0, 0.28,
            font_name='IBM Plex Mono', font_size=10, bold=True, color='#221D17')
add_textbox(slide, '\u20aa34,000 / year', 0.65, ly + 0.42, 3.65, 0.42,
            font_name='Newsreader', font_size=26, color='#B23B26', align=PP_ALIGN.CENTER)
add_textbox(slide, 'in analyst time reclaimed', 0.65, ly + 0.88, 3.65, 0.25,
            font_name='Libre Franklin', font_size=12, color='#5C5349', align=PP_ALIGN.CENTER)

# MIDDLE — dark multiplier
add_rect(slide, 4.52, 1.85, 4.4, 5.25, fill_color='#221D17')
add_textbox(slide, 'THE MULTIPLIER \u2014 TWO EXAMPLES', 4.68, 1.95, 4.1, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#CB5238')

add_textbox(slide, 'Example 1 \u00b7 Contract violation \u00b7 Galil Dairy',
            4.68, 2.22, 4.1, 0.28,
            font_name='Libre Franklin', font_size=12, bold=True, color='#F2ECE0')
ex1 = [
    ('Quoted over contract cap:', '+\u20aa1.30 / kg'),
    ('Volume:', '88 tons / yr'),
    ('Tool caught it \u2014 one meeting:', '\u20aa114,400'),
]
my = 2.55
for label, val in ex1:
    add_textbox(slide, label, 4.68, my, 2.5, 0.28,
                font_name='Libre Franklin', font_size=11, color='#CBC2B2')
    add_textbox(slide, val, 7.1, my, 1.65, 0.28,
                font_name='Libre Franklin', font_size=11, bold=True, color='#F2ECE0',
                align=PP_ALIGN.RIGHT)
    my += 0.35

add_textbox(slide, 'Example 2 \u00b7 Market intelligence \u00b7 Cocoa',
            4.68, my + 0.1, 4.1, 0.28,
            font_name='Libre Franklin', font_size=12, bold=True, color='#F2ECE0')
ex2 = [
    ('Cocoa spot fell:', '\u221229% in 6 months'),
    ('Supplier quote vs contract:', '$3,050 vs $2,650/ton'),
    ('Held contract price \u00b7 35 tons:', '~\u20aa52,000'),
]
my = my + 0.43
for label, val in ex2:
    add_textbox(slide, label, 4.68, my, 2.5, 0.28,
                font_name='Libre Franklin', font_size=11, color='#CBC2B2')
    add_textbox(slide, val, 7.1, my, 1.65, 0.28,
                font_name='Libre Franklin', font_size=11, bold=True, color='#F2ECE0',
                align=PP_ALIGN.RIGHT)
    my += 0.35

add_rect(slide, 4.68, my + 0.05, 3.9, 0.02, fill_color='#3A332A')
add_textbox(slide, '\u20aa114,400 + \u20aa52,000', 4.68, my + 0.12, 3.9, 0.35,
            font_name='Newsreader', font_size=22, color='#CB5238', align=PP_ALIGN.CENTER)
add_textbox(slide, 'Two catches in a 3-day pilot. Tool cost: ~$5/mo.',
            4.68, my + 0.5, 3.9, 0.35,
            font_name='Libre Franklin', font_size=11, color='#A99E8C', align=PP_ALIGN.CENTER,
            wrap=True)

# RIGHT — what to tell CFO
add_rect(slide, 9.0, 1.85, 4.35, 5.25, fill_color='#FBF8F2', border_color='#DDD4C6')
add_textbox(slide, 'WHAT TO TELL THE CFO AT MONTH 3', 9.15, 1.95, 4.0, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#94897A')

cfo_points = [
    'Lead with \u20aa34k/year in analyst time reclaimed.',
    "Don't lead with speculative negotiation savings.",
    'Add: correction rate trending down \u2014 trust improving.',
    'Add: open issues caught before meetings \u2014 count them.',
    'The \u20aa114k multiplier is the kicker \u2014 present it as "one example already found in pilot."',
]
ry = 2.25
for pt in cfo_points:
    add_textbox(slide, '\u2022  ' + pt, 9.15, ry, 4.0, 0.65,
                font_name='Libre Franklin', font_size=12, color='#4A443B', wrap=True)
    ry += 0.72

footer(slide, 18)
set_notes(slide,
    "Conservative anchor first, then two multipliers. "
    "Lead with 34,000 shekels a year in reclaimed analyst time — measurable, auditable, defensible. "
    "Then two concrete examples from the pilot: the Galil Dairy price-cap violation at 114,400 shekels, "
    "and the cocoa benchmark catch at 52,000 shekels in one quarter. "
    "Present both as examples already found, not projections. "
    "If asked about the labor cost number: it uses the manager's own prep time estimate of 1.5 to 3 hours, "
    "midpoint 2.25, across 7 meetings per month. "
    "API cost is approximately $5 per month — it is essentially zero in the ROI calculation.")


# ===========================================================================
# SLIDE 20 — What's Next (dark)
# ===========================================================================

slide = prs.slides.add_slide(blank_layout)
add_slide_bg(slide, '#1F1A14')

# Header
add_textbox(slide, '19 \u2014 WHAT\u2019S NEXT', 0.5, 0.45, 12, 0.3,
            font_name='IBM Plex Mono', font_size=9, color='#CB5238')
add_textbox(slide, 'What I\u2019d prioritise in Phase 2 \u2014 and what I left out',
            0.5, 0.82, 12.5, 1.0,
            font_name='Newsreader', font_size=32, color='#F6F1E8', wrap=True)

# LEFT — priorities
add_textbox(slide, 'PRIORITIES FOR PHASE 2', 0.5, 1.88, 6.0, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#CB5238')

priorities = [
    ('1', 'Live email connector \u2014 Microsoft Graph',
     'Real supplier threads, no synthetic data. Scheduled pull, triggered before each meeting.'),
    ('2', 'SharePoint contract connector',
     'Always the active version, never a stale copy. OCR pipeline for scanned PDFs.'),
    ('3', 'Persist correction log \u2192 prompt-tuning cycle',
     'Monthly review: which fields and suppliers have high error rate? Update the extraction prompt accordingly.'),
    ('4', 'Supplier segmentation',
     'Flag high-risk suppliers by relationship-health score. Auto-generate prep for every meeting, not just on demand.'),
]

py = 2.18
for num, title, desc in priorities:
    add_textbox(slide, num, 0.5, py, 0.4, 0.28,
                font_name='Newsreader', font_size=22, color='#CB5238')
    add_textbox(slide, title, 0.95, py, 5.5, 0.28,
                font_name='Libre Franklin', font_size=14, bold=True, color='#F2ECE0', wrap=True)
    add_textbox(slide, desc, 0.95, py + 0.32, 5.5, 0.5,
                font_name='Libre Franklin', font_size=12, color='#9A9080', wrap=True)
    py += 1.02

# RIGHT — out of scope
add_textbox(slide, 'DELIBERATELY OUT OF SCOPE', 7.1, 1.88, 5.8, 0.25,
            font_name='IBM Plex Mono', font_size=9, color='#4A7360')

out_of_scope = [
    'Negotiation strategy advice \u2014 Human judgment stays human.',
    'Full contract lifecycle management \u2014 Different problem, different tool.',
    'Live MCP connector in prototype \u2014 Demo-day network risk \u2014 documented as production path.',
    'Vector store in prototype \u2014 Unjustified at 8 suppliers / 111 emails \u2014 right call at real scale.',
    'Teams / Slack integration \u2014 Phase 3, after trust is established.',
]
oy = 2.18
for item in out_of_scope:
    add_textbox(slide, '\u00d7  ' + item, 7.1, oy, 6.0, 0.62,
                font_name='Libre Franklin', font_size=13, color='#706558', wrap=True)
    oy += 0.72

footer(slide, 19, dark=True)
set_notes(slide,
    "What I would prioritise in Phase 2 — and what I deliberately left out. "
    "Live connectors first: Microsoft Graph for real emails, SharePoint for always-the-active-contract. "
    "Then the correction-to-prompt-tuning loop: monthly review of which fields and suppliers have high error rate, "
    "update the extraction prompt. "
    "Then supplier segmentation: auto-generate prep for every meeting, not just on demand. "
    "The out-of-scope column matters as much as the priorities "
    "— it shows judgment about what not to build. "
    "Teams integration is Phase 3, after trust is established. "
    "Negotiation strategy advice stays human.")


# ===========================================================================
# Save
# ===========================================================================

output_path = os.path.join(os.path.dirname(__file__), 'Strauss FDE Deck.pptx')
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
